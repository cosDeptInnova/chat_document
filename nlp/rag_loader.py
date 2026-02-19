import os
import json
import logging
import hashlib
import subprocess
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any

import openpyxl
import docx2txt
from pptx import Presentation
from pdfminer.high_level import extract_text as extract_text_from_pdf

logger = logging.getLogger("cosmos_nlp_v3.loader")

# --- Forzar/descubrir providers de ONNXRuntime (GPU/CPU) ---
try:
    import onnxruntime as _ort
except Exception:
    _ort = None  # el entorno decidirá

def _choose_ort_providers() -> tuple[list[str], bool]:
    """
    Devuelve (providers, use_cuda).

    - Respeta DOCLING_OCR_PROVIDERS / RAPIDOCR_PROVIDERS si están definidos.
    - Si no hay ENV, por defecto usa CPU para evitar ocupar VRAM del servidor (multiusuario).
    - Para habilitar CUDA explícitamente: DOCLING_OCR_USE_CUDA=1 o RAG_OCR_USE_CUDA=1.
    """
    # 1) Si el usuario fija providers por ENV, respétalos
    env_providers = (
        os.getenv("DOCLING_OCR_PROVIDERS")
        or os.getenv("RAPIDOCR_PROVIDERS")
        or ""
    ).strip()
    if env_providers:
        providers = [p.strip() for p in env_providers.split(",") if p.strip()]
        use_cuda = any(p.lower().startswith("cuda") for p in providers)
        return providers, use_cuda

    # 2) Flag explícito para permitir GPU (si no está, evitamos CUDA para no llenar VRAM)
    use_cuda_flag = (os.getenv("DOCLING_OCR_USE_CUDA") or os.getenv("RAG_OCR_USE_CUDA") or "0").strip().lower()
    allow_cuda = use_cuda_flag in ("1", "true", "yes", "on")

    # 3) Autodetección (si onnxruntime está importable)
    if _ort is not None:
        try:
            avail = list(_ort.get_available_providers())
        except Exception:
            avail = []
    else:
        avail = []

    # 4) Política: CPU por defecto (multiusuario/VRAM estable)
    if allow_cuda and "CUDAExecutionProvider" in avail:
        return ["CUDAExecutionProvider", "CPUExecutionProvider"], True

    # DirectML opcional (no suele usar VRAM NVIDIA, pero puede existir en hosts raros)
    if "DmlExecutionProvider" in avail:
        return ["DmlExecutionProvider", "CPUExecutionProvider"], False

    return ["CPUExecutionProvider"], False



def clean_text_impl(text: Any) -> str:
    """
    Elimina líneas en blanco, igual que el método original clean_text de DocumentIndexer.
    """
    return "\n".join([line for line in text.splitlines() if line.strip()]) if isinstance(text, str) else ""


def load_txt_impl(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def load_docx_impl(path: str) -> str:
    return docx2txt.process(path)


def load_pdf_impl(indexer, path: str) -> str:
    """
    Extrae texto con pdfminer; si es escaso, hace OCR (spa+eng) con pdf2image+pytesseract.
    Mantiene la lógica del método original load_pdf.
    """
    try:
        txt = extract_text_from_pdf(path) or ""
    except Exception:
        txt = ""
    cleaned = clean_text_impl(txt)
    if len(cleaned) >= 300:
        return cleaned

    # OCR fallback
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(path, dpi=300)
        ocr_txt = []
        for im in images[:50]:  # límite de páginas por seguridad
            ocr_txt.append(pytesseract.image_to_string(im, lang="spa+eng"))
        return clean_text_impl("\n".join(ocr_txt))
    except Exception as e:
        logger.warning(f"(pdf_ocr) Fallback OCR falló: {e}")
        return cleaned


def load_pptx_impl(path: str) -> str:
    prs = Presentation(path)
    return "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


# ================================================================
# Helpers de perfilado de tablas Excel
# ================================================================

def build_table_profile(table_id: str, df: Any) -> Dict[str, Any]:
    """
    Construye un perfil estructurado de una tabla (DataFrame):
    - nº filas/columnas
    - por columna: dtype, nulls, n_unique, top-K valores, stats numéricas básicas.
    """
    profile: Dict[str, Any] = {
        "table_id": table_id,
        "n_rows": int(len(df)),
        "n_cols": int(len(df.columns)),
        "columns": {},
    }

    topk = int(os.getenv("EXCEL_PROFILE_TOPK", "10"))

    for col in df.columns:
        series = df[col]
        col_info: Dict[str, Any] = {
            "dtype": str(series.dtype),
            "non_null": int(series.notna().sum()),
            "nulls": int(series.isna().sum()),
        }

        # cardinalidad
        try:
            col_info["n_unique"] = int(series.nunique(dropna=True))
        except Exception:
            col_info["n_unique"] = None

        # top-K valores
        try:
            vc = series.value_counts(dropna=True).head(topk)
            col_info["top_values"] = [
                {"value": str(idx), "count": int(cnt)} for idx, cnt in vc.items()
            ]
        except Exception:
            col_info["top_values"] = []

        # stats numéricas básicas
        try:
            # intentamos convertir a float de forma robusta
            s_num = series.dropna()
            s_num = s_num.astype(float)
            if not s_num.empty:
                col_info["min"] = float(s_num.min())
                col_info["max"] = float(s_num.max())
                col_info["mean"] = float(s_num.mean())
                col_info["sum"] = float(s_num.sum())
        except Exception:
            pass

        profile["columns"][str(col)] = col_info

    return profile


def profile_to_text(table_id: str, profile: Dict[str, Any]) -> str:
    """
    Convierte el perfil de tabla en un texto resumido apto para embeddings.
    """
    n_rows = profile.get("n_rows")
    n_cols = profile.get("n_cols")
    header = f"Resumen de la tabla {table_id}: {n_rows} filas, {n_cols} columnas."

    col_summaries: List[str] = []
    for col_name, col_info in (profile.get("columns") or {}).items():
        top_vals = col_info.get("top_values") or []
        if not top_vals:
            continue
        tops = ", ".join(f"{tv['value']} ({tv['count']})" for tv in top_vals[:3])
        col_summaries.append(f"{col_name}: {tops}")

    if col_summaries:
        return header + " Valores más frecuentes por columna: " + "; ".join(col_summaries)
    return header


# ================================================================
# Backends de parsing enriquecido: Docling y MinerU
# ================================================================
DOC_STRUCT_BACKEND = os.getenv(
    "RAG_STRUCTURED_BACKEND",
    os.getenv("DOC_STRUCTURED_BACKEND", "docling"),
).lower()

# --------------- Docling ---------------
try:
    from docling.document_converter import DocumentConverter
    from docling.chunking import HybridChunker

    _HAS_DOCLING = True
except Exception:  # ImportError o fallo de inicialización perezosa
    DocumentConverter = None  # type: ignore
    HybridChunker = None  # type: ignore
    _HAS_DOCLING = False

_DOCLING_CONVERTER: Optional[Any] = None
_DOCLING_CHUNKER: Optional[Any] = None


def _get_docling() -> Tuple[Optional[Any], Optional[Any]]:
    """
    Inicializa (una vez) DocumentConverter + HybridChunker de Docling.

    - Usa providers de _choose_ort_providers() (CPU por defecto).
    - Intenta pasar rapidocr_kwargs si la versión lo soporta.
    - Si no, fuerza ENV compatibles con versiones antiguas.
    """
    global _DOCLING_CONVERTER, _DOCLING_CHUNKER

    if not _HAS_DOCLING:
        return None, None

    if _DOCLING_CONVERTER is not None and _DOCLING_CHUNKER is not None:
        return _DOCLING_CONVERTER, _DOCLING_CHUNKER

    try:
        providers, use_cuda = _choose_ort_providers()
        ocr_engine = os.getenv("DOCLING_OCR_ENGINE", "rapidocr")

        opts = None
        try:
            from docling.pipeline.standard_pdf_pipeline import StandardPdfPipelineOptions
            try:
                opts = StandardPdfPipelineOptions(
                    ocr_engine=ocr_engine,
                    rapidocr_kwargs={
                        "use_cuda": use_cuda,
                        "providers": providers,
                    },
                )
                logger.info(
                    "[docling] ocr_engine=%s con rapidocr_kwargs (providers=%s, use_cuda=%s)",
                    ocr_engine, providers, use_cuda,
                )
            except TypeError:
                opts = StandardPdfPipelineOptions(ocr_engine=ocr_engine)
                raise ImportError("StandardPdfPipelineOptions sin rapidocr_kwargs")
        except Exception:
            # Docling “antiguo” → forzamos ENV legacy
            os.environ["DOCLING_OCR_ENGINE"] = ocr_engine
            os.environ["DOCLING_ACCELERATOR"] = "cuda" if use_cuda else "cpu"

            os.environ["RAPIDOCR_DEVICE"] = "cuda" if use_cuda else "cpu"
            os.environ["RAPIDOCR_ENGINE_NAME"] = "onnxruntime"
            os.environ["RAPIDOCR_ONNX_PROVIDERS"] = ",".join(providers)
            os.environ["RAPIDOCR_PROVIDERS"] = ",".join(providers)

            logger.info(
                "[docling] Sin rapidocr_kwargs; forzando ENV "
                "DOCLING_OCR_ENGINE=%s DOCLING_ACCELERATOR=%s RAPIDOCR_DEVICE=%s "
                "RAPIDOCR_ENGINE_NAME=%s RAPIDOCR_ONNX_PROVIDERS=%s",
                os.environ.get("DOCLING_OCR_ENGINE"),
                os.environ.get("DOCLING_ACCELERATOR"),
                os.environ.get("RAPIDOCR_DEVICE"),
                os.environ.get("RAPIDOCR_ENGINE_NAME"),
                os.environ.get("RAPIDOCR_ONNX_PROVIDERS"),
            )
            opts = None

        if opts is not None:
            _DOCLING_CONVERTER = DocumentConverter(pipeline_options=opts)
        else:
            _DOCLING_CONVERTER = DocumentConverter()

        _DOCLING_CHUNKER = HybridChunker()

        try:
            import onnxruntime as ort
            logger.info("[docling] ORT providers disponibles: %s", ort.get_available_providers())
        except Exception:
            pass

        logger.info("[docling] DocumentConverter + HybridChunker inicializados correctamente.")
    except Exception as e:
        logger.warning(f"[docling] No se pudo inicializar Docling, se usará el pipeline clásico: {e}")
        _DOCLING_CONVERTER = None
        _DOCLING_CHUNKER = None

    return _DOCLING_CONVERTER, _DOCLING_CHUNKER


def _load_with_docling(indexer, path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Usa Docling + HybridChunker para convertir el fichero en chunks enriquecidos.
    Añade limitador de concurrencia para evitar picos de memoria.

    ENV:
      DOCLING_MAX_CONCURRENCY (default=1)
    """
    import threading

    converter, chunker = _get_docling()
    if converter is None or chunker is None:
        return []

    # Semáforo global (una vez) para limitar conversions concurrentes
    g = globals()
    if "_DOCLING_SEM" not in g:
        try:
            cap = max(1, int(os.getenv("DOCLING_MAX_CONCURRENCY", "1")))
        except Exception:
            cap = 1
        g["_DOCLING_SEM"] = threading.BoundedSemaphore(value=cap)

    sem = g["_DOCLING_SEM"]

    try:
        sem.acquire()
        try:
            result = converter.convert(source=path)
            dl_doc = result.document
        finally:
            # el chunking también puede ser pesado; lo mantenemos dentro del cupo
            pass
    except Exception as e:  # pragma: no cover
        try:
            sem.release()
        except Exception:
            pass
        logger.warning(f"[docling] Error convirtiendo {path}: {e}")
        return []

    try:
        chunk_iter = chunker.chunk(dl_doc=dl_doc)
    except Exception as e:  # pragma: no cover
        try:
            sem.release()
        except Exception:
            pass
        logger.warning(f"[docling] Error chunking {path}: {e}")
        return []

    chunks_out: List[Tuple[str, Dict[str, Any]]] = []
    base = os.path.basename(path)

    try:
        for idx, chunk in enumerate(chunk_iter):
            try:
                enriched_text = chunker.contextualize(chunk=chunk)
            except Exception:  # pragma: no cover
                enriched_text = getattr(chunk, "text", "")

            enriched_text = clean_text_impl(enriched_text)
            if not enriched_text:
                continue

            meta: Dict[str, Any] = {
                "source_path": path,
                "source_file_name": base,
                "backend": "docling",
                "chunk_index": idx,
            }

            if hasattr(chunk, "meta") and getattr(chunk, "meta") is not None:
                try:
                    docling_meta = chunk.meta.export_json_dict()  # type: ignore[attr-defined]
                    if isinstance(docling_meta, dict):
                        meta.update(docling_meta)

                        tags: List[str] = []
                        for k, v in docling_meta.items():
                            kl = str(k).lower()
                            if any(t in kl for t in ("heading", "title", "section")):
                                if isinstance(v, str):
                                    if v.strip():
                                        tags.append(v.strip())
                                elif isinstance(v, (list, tuple)):
                                    for vv in v:
                                        sv = str(vv).strip()
                                        if sv:
                                            tags.append(sv)
                        if tags:
                            meta.setdefault("tags", tags)
                except Exception as e:  # pragma: no cover
                    logger.debug(f"[docling] No se pudo serializar metadata para {path}: {e}")

            chunks_out.append((enriched_text, meta))
    finally:
        try:
            sem.release()
        except Exception:
            pass

    return chunks_out


# --------------- MinerU (opcional, vía CLI) ---------------
MINERU_CLI = os.getenv("MINERU_CLI", "mineru")
MINERU_OUTPUT_DIR = os.getenv("MINERU_OUTPUT_DIR", "./_mineru_cache")
_MINERU_AVAILABLE: Optional[bool] = None


def _load_with_mineru(path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Usa MinerU vía CLI para PDFs.

    - Ejecuta: `mineru -p <path> -o <cache_dir>`
    - Intenta leer `{basename}_content_list.json` para construir chunks con metadata.
    - Si falla, hace fallback a `<basename>.md` y lo trocea por párrafos.

    Requiere tener MinerU instalado y accesible en PATH como `mineru` o via MINERU_CLI.
    """
    global _MINERU_AVAILABLE

    if _MINERU_AVAILABLE is False:
        return []

    # Comprobación perezosa de disponibilidad del binario
    if _MINERU_AVAILABLE is None:
        try:
            from shutil import which

            if which(MINERU_CLI) is None:
                logger.warning(
                    "[mineru] Binario '%s' no encontrado en PATH; desactivando backend MinerU.",
                    MINERU_CLI,
                )
                _MINERU_AVAILABLE = False
                return []
            _MINERU_AVAILABLE = True
        except Exception:
            _MINERU_AVAILABLE = False
            return []

    base_dir = Path(MINERU_OUTPUT_DIR).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)

    file_hash = hashlib.sha1(path.encode("utf-8")).hexdigest()
    out_dir = base_dir / file_hash
    out_dir.mkdir(parents=True, exist_ok=True)

    stem = Path(path).stem
    md_path = out_dir / f"{stem}.md"
    cl_path = out_dir / f"{stem}_content_list.json"

    # Ejecutar MinerU sólo si no tenemos outputs en cache
    if not (md_path.exists() or cl_path.exists()):
        cmd = [MINERU_CLI, "-p", path, "-o", str(out_dir)]
        try:
            res = subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
            logger.info("[mineru] Ejecutado correctamente para %s: %s", path, res.args)
        except Exception as e:  # pragma: no cover
            logger.warning(f"[mineru] Falló la ejecución para {path}: {e}")
            return []

    chunks_out: List[Tuple[str, Dict[str, Any]]] = []

    # 1) Intentar structured JSON (content_list.json)
    if cl_path.exists():
        try:
            with open(cl_path, "r", encoding="utf-8") as f:
                content_list = json.load(f)

            if isinstance(content_list, list):
                for idx, block in enumerate(content_list):
                    if not isinstance(block, dict):
                        continue

                    text = ""
                    for key in ("markdown", "md", "text", "content"):
                        v = block.get(key)
                        if isinstance(v, str) and v.strip():
                            text = v
                            break
                    text = clean_text_impl(text)
                    if not text:
                        continue

                    meta: Dict[str, Any] = {
                        "source_path": path,
                        "source_file_name": stem,
                        "backend": "mineru",
                        "block_index": idx,
                        "block_type": block.get("type") or block.get("block_type"),
                        "page": block.get("page") or block.get("page_number"),
                    }
                    if "bbox" in block:
                        meta["bbox"] = block["bbox"]
                    chunks_out.append((text, meta))
        except Exception as e:  # pragma: no cover
            logger.warning(f"[mineru] Error leyendo content_list.json para {path}: {e}")

    # 2) Fallback: sólo Markdown → fragmentar por párrafos
    if not chunks_out and md_path.exists():
        try:
            with open(md_path, "r", encoding="utf-8") as f:
                md_text = f.read()
            # Troceo muy simple; el chunking semántico adicional lo puede hacer
            # fragment_text_semantically_impl si lo deseamos.
            paragraphs = [p.strip() for p in md_text.split("\n\n") if p.strip()]
            for idx, para in enumerate(paragraphs):
                chunks_out.append(
                    (
                        para,
                        {
                            "source_path": path,
                            "source_file_name": stem,
                            "backend": "mineru",
                            "block_index": idx,
                        },
                    )
                )
        except Exception as e:  # pragma: no cover
            logger.warning(f"[mineru] Error leyendo markdown para {path}: {e}")

    return chunks_out


def load_and_fragment_files_impl(indexer, files: List[str]):
    """
    Versión evolucionada de DocumentIndexer.load_and_fragment_files.

    Mejoras clave (sin romper el pipeline):
    - Excel: read_only=True + wb.close() + no listas enormes de índices por hoja.
    - Mantiene Docling/MinerU y pipeline clásico como estaban.
    """
    from tqdm import tqdm

    documents: List[str] = []
    doc_names: List[str] = []

    for path in tqdm(files, desc="Procesando archivos"):
        print(f"(load) Procesando archivo: {path}", flush=True)
        ext = os.path.splitext(path.lower())[1]

        try:
            # ------------------------------------------------------------
            # 1) Pipeline especializado para Excel (enriquecido)
            # ------------------------------------------------------------
            if ext == ".xlsx":
                wb = None
                try:
                    # read_only reduce RAM brutal en Excels grandes
                    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
                    basename = os.path.basename(path)

                    MACRO_CHARS = int(os.getenv("EXCEL_MACRO_CHUNK_CHARS", "1500"))
                    EXCEL_ROW_JOIN = os.getenv("EXCEL_ROW_JOIN", " | ")
                    EXCEL_HEADER_PREFIX = os.getenv("EXCEL_HEADER_PREFIX", "")
                    EXCEL_INCLUDE_SHEET_MACRO = os.getenv("EXCEL_INCLUDE_SHEET_MACRO", "1") == "1"
                    EXCEL_TABLE_CACHE_DIR = os.getenv("EXCEL_TABLE_CACHE_DIR", "./_excel_tables")

                    try:
                        import pandas as pd  # type: ignore
                        pd_available = True
                    except Exception:
                        pd_available = False
                        pd = None  # type: ignore

                    table_cache_root: Optional[Path] = None
                    if pd_available:
                        table_cache_root = Path(EXCEL_TABLE_CACHE_DIR).resolve()
                        table_cache_root.mkdir(parents=True, exist_ok=True)

                    try:
                        _aliases_raw = indexer._standard_field_aliases()
                    except Exception:
                        _aliases_raw = {}

                    alias_map_norm: Dict[str, str] = {}
                    if _aliases_raw:
                        try:
                            alias_map_norm = {indexer._canon_key(k): v for k, v in _aliases_raw.items()}
                        except Exception:
                            alias_map_norm = {}

                    def _apple_class_fix(row_kv: Dict[str, str]) -> Dict[str, str]:
                        kv = dict(row_kv or {})
                        canon = {indexer._canon(k): str(v) for k, v in kv.items()}
                        cls = canon.get("class") or canon.get("tipo") or ""

                        model = canon.get("model") or canon.get("modelo") or ""
                        brand = canon.get("brand") or canon.get("marca") or ""
                        joined = f"{brand} {model}".upper()

                        def set_class(v):
                            kv["class"] = v
                            return kv

                        if cls:
                            c = cls.lower()
                            if any(t in c for t in ["imac", "ipad", "macbook"]):
                                return kv
                        if any(a in joined for a in ["A1419", "A1418"]) or "IMAC" in joined:
                            return set_class("imac")
                        if any(a in joined for a in ["PRO 16", "PRO 14", "MACBOOK"]):
                            return set_class("macbook_pro")
                        if "IPAD" in joined:
                            return set_class("ipad")
                        return kv

                    def _canonicalize_row_meta(row_kv: Dict[str, str]) -> Dict[str, Any]:
                        canon_kv: Dict[str, str] = {}
                        numeric_fields: Dict[str, float] = {}

                        for raw_key, val in (row_kv or {}).items():
                            if not raw_key:
                                continue
                            try:
                                ck_norm = indexer._canon_key(str(raw_key))
                            except Exception:
                                continue
                            canonical_name = alias_map_norm.get(ck_norm)
                            if not canonical_name:
                                continue
                            canon_kv[canonical_name] = val
                            if canonical_name in ("quantity", "price"):
                                try:
                                    num = indexer._parse_number_es(val)
                                    if num is not None:
                                        numeric_fields[canonical_name] = num
                                except Exception:
                                    pass

                        meta_extra: Dict[str, Any] = {}
                        for key in ("location", "class", "brand", "model", "asset_id", "serial", "asset_type"):
                            if key in canon_kv:
                                meta_extra[key] = canon_kv[key]

                        if "quantity" in numeric_fields:
                            meta_extra["quantity"] = numeric_fields["quantity"]
                        if "price" in numeric_fields:
                            meta_extra["price"] = numeric_fields["price"]

                        if canon_kv:
                            meta_extra["row_kv_canon"] = canon_kv

                        return meta_extra

                    for sheet in wb.worksheets:
                        print(f"(load)     Worksheet: {sheet.title}", flush=True)
                        try:
                            rows_raw = list(sheet.iter_rows(values_only=True))
                        except Exception:
                            logger.warning(f"(load)   Worksheet inválida/dañada: {sheet.title}")
                            continue

                        rows_str: List[List[str]] = []
                        for row in rows_raw:
                            rows_str.append([(str(c) if c is not None else "").strip() for c in row])

                        hdr_row_index: Optional[int] = indexer._detect_header_row(rows_str)
                        headers: List[str] = []
                        start_idx = 0
                        if hdr_row_index is not None:
                            headers = [h.strip() for h in rows_str[hdr_row_index]]
                            start_idx = hdr_row_index + 1

                        #  en vez de acumular indices (lista gigante), guardamos rango de metas por hoja
                        sheet_meta_start = len(indexer._metas)

                        sheet_rows_for_df: List[Dict[str, Any]] = []

                        if EXCEL_INCLUDE_SHEET_MACRO:
                            macro_lines: List[str] = []
                            for r in rows_str[start_idx:]:
                                joined = EXCEL_ROW_JOIN.join([c for c in r if c])
                                if joined:
                                    macro_lines.append(joined)

                            chunk_buf: List[str] = []
                            cur = 0
                            chunk_idx = 1

                            def flush_macro():
                                nonlocal chunk_buf, cur, chunk_idx
                                if not chunk_buf:
                                    return
                                macro_text = "\n".join(chunk_buf)
                                macro_text = "\n".join([ln for ln in macro_text.splitlines() if ln.strip()])
                                if macro_text:
                                    documents.append(macro_text)
                                    doc_names.append(f"{basename}::{sheet.title}::macro#{chunk_idx}")
                                    indexer._metas.append({
                                        "sheet": sheet.title,
                                        "row_idx": None,
                                        "headers": headers,
                                        "source_path": path,
                                        "backend": "excel_macro",
                                    })
                                chunk_buf = []
                                cur = 0
                                chunk_idx += 1

                            for line in macro_lines:
                                if cur + len(line) + 1 > MACRO_CHARS:
                                    flush_macro()
                                chunk_buf.append(line)
                                cur += len(line) + 1
                            flush_macro()

                        for ridx, row in enumerate(rows_str[start_idx:], start=start_idx + 1):
                            row_kv: Dict[str, str] = {}
                            pairs: List[str] = []
                            for i, v in enumerate(row):
                                if not v:
                                    continue
                                if headers and i < len(headers) and headers[i].strip():
                                    h = headers[i].strip()
                                    row_kv[h] = v
                                    pairs.append(f"{h}={v}" if not EXCEL_HEADER_PREFIX else f"{h}: {v}")
                                else:
                                    synth_key = f"col_{i+1}"
                                    row_kv[synth_key] = v
                                    pairs.append(v)

                            if row_kv:
                                row_kv = _apple_class_fix(row_kv)

                            row_text = EXCEL_ROW_JOIN.join(pairs).strip()
                            if not row_text:
                                continue

                            documents.append(row_text)
                            doc_names.append(basename)

                            meta = {
                                "sheet": sheet.title,
                                "row_idx": ridx,
                                "headers": headers,
                                "row_kv": row_kv,
                                "source_path": path,
                                "backend": "excel_row",
                            }
                            extra_meta = _canonicalize_row_meta(row_kv)
                            meta.update(extra_meta)

                            indexer._metas.append(meta)

                            if pd_available and row_kv:
                                row_for_df: Dict[str, Any] = dict(row_kv)
                                for k in ("location", "class", "brand", "model",
                                          "asset_id", "serial", "asset_type",
                                          "quantity", "price"):
                                    if k in extra_meta:
                                        row_for_df[k] = extra_meta[k]
                                sheet_rows_for_df.append(row_for_df)

                        # Cache tabular + perfil por hoja
                        if pd_available and table_cache_root is not None and sheet_rows_for_df:
                            try:
                                table_id = f"{basename}::{sheet.title}"
                                slug = table_id.replace("::", "__")
                                table_cache_path = table_cache_root / f"{slug}.parquet"

                                df = pd.DataFrame(sheet_rows_for_df)
                                df.to_parquet(table_cache_path)

                                profile = build_table_profile(table_id, df)
                                profile_path = table_cache_root / f"{slug}_profile.json"
                                profile_path.write_text(
                                    json.dumps(profile, ensure_ascii=False, indent=2),
                                    encoding="utf-8",
                                )

                                #  Propagar table_id/cache_path a las metas de la hoja por rango (sin lista gigante)
                                sheet_meta_end = len(indexer._metas)
                                for mi in range(sheet_meta_start, sheet_meta_end):
                                    try:
                                        m = indexer._metas[mi]
                                        m["table_id"] = table_id
                                        m["table_cache_path"] = str(table_cache_path)
                                    except Exception:
                                        pass

                                profile_text = profile_to_text(table_id, profile)
                                documents.append(profile_text)
                                doc_names.append(f"{basename}::{sheet.title}::profile")
                                indexer._metas.append({
                                    "sheet": sheet.title,
                                    "row_idx": None,
                                    "headers": headers,
                                    "source_path": path,
                                    "backend": "excel_profile",
                                    "table_id": table_id,
                                    "table_cache_path": str(table_cache_path),
                                    "profile_path": str(profile_path),
                                })

                                print(f"(load)   Excel: cache/parquet y perfil generados para {table_id}", flush=True)
                            except Exception as e:
                                logger.warning(f"(load)   No se pudo generar cache/perfil para Excel {basename} hoja {sheet.title}: {e}")

                    continue  # siguiente archivo
                finally:
                    try:
                        if wb is not None:
                            wb.close()
                    except Exception:
                        pass

            # ------------------------------------------------------------
            # 2) Backends estructurados Docling / MinerU
            # ------------------------------------------------------------
            used_structured = False
            struct_exts = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".html", ".htm", ".md"}

            if ext in struct_exts:
                backend = DOC_STRUCT_BACKEND

                if backend in ("docling", "auto"):
                    dl_chunks = _load_with_docling(indexer, path)
                    if dl_chunks:
                        used_structured = True
                        for text, meta in dl_chunks:
                            documents.append(text)
                            doc_names.append(os.path.basename(path))
                            indexer._metas.append(meta)
                        print(f"(load)   Docling: {len(dl_chunks)} fragmentos generados para {path}", flush=True)

                if not used_structured and backend == "mineru" and ext == ".pdf":
                    mineru_chunks = _load_with_mineru(path)
                    if mineru_chunks:
                        used_structured = True
                        for text, meta in mineru_chunks:
                            documents.append(text)
                            doc_names.append(os.path.basename(path))
                            indexer._metas.append(meta)
                        print(f"(load)   MinerU: {len(mineru_chunks)} fragmentos generados para {path}", flush=True)

                if used_structured:
                    continue

            # ------------------------------------------------------------
            # 3) Pipeline clásico (fallback)
            # ------------------------------------------------------------
            if ext == ".txt":
                print(f"(load)   Texto plano detectado", flush=True)
                text = load_txt_impl(path)
            elif ext == ".docx":
                print(f"(load)   Word detectado (pipeline clásico)", flush=True)
                text = load_docx_impl(path)
            elif ext == ".pdf":
                print(f"(load)   PDF detectado (pipeline clásico)", flush=True)
                text = load_pdf_impl(indexer, path)
            elif ext == ".pptx":
                print(f"(load)   PowerPoint detectado (pipeline clásico)", flush=True)
                text = load_pptx_impl(path)
            else:
                print(f"(load)   Extensión no válida o no soportada: {ext}", flush=True)
                continue

            text = clean_text_impl(text)
            if not text:
                continue

            est_tokens = max(1, len(text) // 4)
            size = indexer.dynamic_fragment_size(est_tokens)
            print(f"(load)   Fragmentando: est_tokens={est_tokens}, tamaño_fragmento={size}", flush=True)

            fragments = indexer.fragment_text_semantically(text, max_tokens=size, overlap_tokens=33)
            for frag in fragments:
                documents.append(frag)
                doc_names.append(os.path.basename(path))
                indexer._metas.append({"source_path": path, "backend": "classic"})

        except Exception as e:
            logger.warning(f"Error procesando {path}: {e}")

    print(f"(load) Total fragments: {len(documents)}", flush=True)
    return documents, doc_names, False


_SPACY_NLP = None


def _get_nlp():
    global _SPACY_NLP
    if _SPACY_NLP is None:
        import spacy

        model = os.getenv("SPACY_MODEL_ES", "es_core_news_md")

        # Desactiva componentes pesados que NO son necesarios para .sents (mantén parser/senter)
        disable_env = os.getenv("SPACY_DISABLE", "ner")
        disable = [c.strip() for c in disable_env.split(",") if c.strip()]
        disable = [c for c in disable if c not in ("parser", "senter", "sentencizer")]

        _SPACY_NLP = spacy.load(model, disable=disable)

        # Evita errores por textos grandes
        try:
            _SPACY_NLP.max_length = int(os.getenv("SPACY_MAX_LENGTH", "5000000"))
        except Exception:
            pass

        # Si por configuración el pipeline no produce sents, añadimos sentencizer
        try:
            has_sents = any(p in _SPACY_NLP.pipe_names for p in ("parser", "senter", "sentencizer"))
            if not has_sents:
                _SPACY_NLP.add_pipe("sentencizer")
        except Exception:
            pass

    return _SPACY_NLP


def dynamic_fragment_size_impl(total_tokens: int) -> int:
    """
    Se mantiene firma y trazas (print) del original.
    """
    print(f"(fragment_size) total_tokens={total_tokens}", flush=True)
    size = 60 if total_tokens < 240 else 120 if total_tokens < 1200 else 180
    print(f"(fragment_size) size={size}", flush=True)
    return size


def fragment_text_semantically_impl(
    text: str,
    max_tokens: int = 100,
    overlap_tokens: int = 33,
) -> List[str]:
    """
    Segmentación por oraciones con refuerzo semántico en UNA sola pasada de spaCy.
    Misma lógica, pero sin construir listas enormes intermedias (streaming).
    """
    import re

    print(
        f"(fragment) Iniciando fragmentación semántica: max_tokens={max_tokens}, overlap={overlap_tokens}",
        flush=True,
    )

    title_re = re.compile(r"^\s*([A-ZÁÉÍÓÚÜÑ0-9][A-ZÁÉÍÓÚÜÑ0-9 \-_]{3,}:?|#{1,6}\s+.+)$")
    bullet_re = re.compile(r"^\s*(?:[-*•\u2022]|\d{1,2}\.)\s+")

    lines = [ln.rstrip() for ln in text.splitlines()]
    blocks: List[str] = []
    buf: List[str] = []
    for ln in lines:
        if title_re.match(ln) and buf:
            blocks.append("\n".join(buf).strip())
            buf = [ln]
        else:
            buf.append(ln)
    if buf:
        blocks.append("\n".join(buf).strip())

    def split_block(btxt: str) -> List[str]:
        parts: List[str] = []
        cur: List[str] = []
        for ln in btxt.splitlines():
            if bullet_re.match(ln) or not ln.strip():
                cur.append(ln)
            else:
                if cur:
                    parts.append("\n".join(cur).strip())
                    cur = []
                parts.append(ln)
        if cur:
            parts.append("\n".join(cur).strip())
        return [p for p in parts if p.strip()]

    segments: List[str] = []
    for b in blocks:
        segments.extend(split_block(b))

    nlp = _get_nlp()
    batch_size = int(os.getenv("NLP_PIPE_BATCH", "32"))

    fragments: List[str] = []
    current: List[str] = []
    counts: List[int] = []
    total = 0

    def _flush_current():
        nonlocal current, counts, total
        frag = " ".join(current).strip()
        if frag:
            fragments.append(frag)

        # overlap (igual que antes)
        overlap: List[str] = []
        o_counts: List[int] = []
        tot = 0
        for s, cc in reversed(list(zip(current, counts))):
            if not s.strip():
                continue
            overlap.insert(0, s)
            o_counts.insert(0, cc)
            tot += cc
            if tot >= overlap_tokens:
                break
        current, counts = overlap, o_counts
        total = sum(o_counts)

    # Streaming spaCy: no guardamos docs ni todas las frases
    for d in nlp.pipe(segments, batch_size=batch_size):
        for s in d.sents:
            st = s.text.strip()
            if not st:
                continue
            tk = len(s)

            current.append(st)
            counts.append(tk)
            total += tk

            if total >= max_tokens:
                _flush_current()

    if current:
        frag = " ".join(current).strip()
        if frag:
            fragments.append(frag)

    print(f"(fragment) Total fragments generados: {len(fragments)}", flush=True)
    return fragments
