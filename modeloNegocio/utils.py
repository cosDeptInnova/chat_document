import os
import json
import uuid
import secrets
import time
import logging
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import asyncio
import httpx
import pandas as pd
import docx2txt
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
from fastapi import Request, HTTPException, status
from jose import JWTError, jwt
from prometheus_client import Counter, Histogram
from sqlalchemy.orm import Session
import os
import socket
import struct
import time
import asyncio
import logging
from typing import Dict, Any, Optional
import hmac
import base64
import hashlib
from urllib.parse import urlparse, urlencode


from config.models import Conversation

# =============================================================================
#  CONFIGURACIÓN GENERAL
# =============================================================================

# Clave JWT y algoritmo: mismos que login / nlp / addUser
SECRET_KEY = os.getenv("COSMOS_SECRET_KEY", "secretkey123")
ALGORITHM = os.getenv("COSMOS_JWT_ALG", "HS256")

# Raíz de almacenamiento compartida por todos los servicios
STORAGE_ROOT = Path(os.getenv("COSMOS_STORAGE_ROOT", "./cosmos_data")).resolve()

# URLs de servicios internos (configurables para Docker o VS Code)
LLM_URL = os.getenv(
    "COSMOS_LLM_URL_RAG",
    "http://localhost:8081/generate_response_with_context_rag",
)
LLAMA_SERVER_URL = os.getenv(
    "COSMOS_LLM_URL_BASE",
    "http://localhost:8081/generate_response",
)
LLM_CONTEXT_CONVERSATION = os.getenv(
    "COSMOS_LLM_URL_CONV",
    "http://localhost:8081/generate_response_conversation",
)
LLM_AI_IMAGE = os.getenv(
    "COSMOS_LLM_URL_IMAGE",
    "http://localhost:8081/generate_response_ai_image",
)
LLM_LAST_DOCUMENT = os.getenv(
    "COSMOS_LLM_URL_LASTDOC",
    "http://localhost:8081/ask_last_document",
)

# OCR / NLP / TalkToDocument
OCR_APP_URL = os.getenv("COSMOS_OCR_URL", "http://localhost:6000")
NLP_APP_URL = os.getenv("COSMOS_NLP_URL", "http://localhost:5000")
TALK_TO_DOCUMENT_URL = os.getenv(
    "COSMOS_TALK_TO_DOC_URL", "http://localhost:5001/chat_document"
)
NLP_UPLOAD_FILE_URL = os.getenv(
    "COSMOS_NLP_UPLOAD_URL", "http://localhost:5000/upload_file"
)

# Roles (IDs numéricos en BD)
ROLE_SUPERVISOR = 1
ROLE_USER = 2

logger = logging.getLogger(__name__)

# =============================================================================
#  MÉTRICAS PROMETHEUS
# =============================================================================

REQUEST_COUNT = Counter(
    "http_requests_total",
    "Total de solicitudes HTTP",
    ["method", "endpoint"],
)

RESPONSE_TIME = Histogram(
    "http_response_time_seconds",
    "Tiempo de respuesta HTTP",
    ["method", "endpoint"],
)

HTTP_REQUEST_DURATION = Histogram(
    "http_request_duration_seconds",
    "Duración de las peticiones HTTP",
    ["method", "endpoint"],
)

# JWT verification
JWT_VERIFY_REQUESTS = Counter(
    "jwt_verify_requests_total",
    "Número de veces que se llamó a verify_token_from_cookie",
)
JWT_MISSING = Counter(
    "jwt_verify_missing_token_total",
    "Veces que no se encontró token en la cookie",
)
JWT_INVALID = Counter(
    "jwt_verify_invalid_token_total",
    "Veces que el token JWT fue inválido o expirado",
)
JWT_VALID = Counter(
    "jwt_verify_valid_token_total",
    "Veces que el token JWT fue verificado correctamente",
)
JWT_VERIFY_DURATION = Histogram(
    "jwt_verify_duration_seconds",
    "Tiempo que tarda la verificación de token JWT",
)

# Redis: sesión
REDIS_SESSION_REQUESTS = Counter(
    "redis_session_requests_total",
    "Total calls to get_session_from_redis",
)
REDIS_SESSION_HITS = Counter(
    "redis_session_hits_total",
    "Redis session cache hits",
)
REDIS_SESSION_MISSES = Counter(
    "redis_session_misses_total",
    "Redis session cache misses",
)
REDIS_SESSION_DURATION = Histogram(
    "redis_session_duration_seconds",
    "Duration of get_session_from_redis",
)

# Redis: conversaciones
CONV_GET_REQUESTS = Counter(
    "redis_conversation_get_requests_total",
    "Total calls to get_conversation_to_redis",
)
CONV_GET_HITS = Counter(
    "redis_conversation_get_hits_total",
    "Redis conversation cache hits",
)
CONV_GET_MISSES = Counter(
    "redis_conversation_get_misses_total",
    "Redis conversation cache misses",
)
CONV_GET_DURATION = Histogram(
    "redis_conversation_get_duration_seconds",
    "Duration of get_conversation_to_redis",
)

CONV_SAVE_REQUESTS = Counter(
    "redis_conversation_save_requests_total",
    "Total calls to save_conversation_to_redis",
)
CONV_SAVE_DURATION = Histogram(
    "redis_conversation_save_duration_seconds",
    "Duration of save_conversation_to_redis",
)

CONV_RESET_REQUESTS = Counter(
    "redis_conversation_reset_requests_total",
    "Total calls to reset_conversation_in_redis",
)
CONV_RESET_DURATION = Histogram(
    "redis_conversation_reset_duration_seconds",
    "Duration of reset_conversation_in_redis",
)

# Upload de ficheros efímeros
UPLOADFILE_REQUESTS = Counter(
    "uploadfile_requests_total",
    "Total calls to create_upload_file",
)
UPLOADFILE_ERRORS = Counter(
    "uploadfile_errors_total",
    "Total failures in create_upload_file",
)
UPLOADFILE_DURATION = Histogram(
    "uploadfile_duration_seconds",
    "Duration of create_upload_file handler",
    ["method", "endpoint"],
)

# LLM queries
llm_query_counter = Counter(
    "query_llm_requests_total",
    "Número total de peticiones a /query/llm",
)

llm_query_requests_by_status_total = Counter(
    "llm_query_requests_by_status_total",
    "Número total de peticiones a /query/llm clasificadas por status",
    ["status"],
)

llm_query_latency = Histogram(
    "llm_query_latency_seconds",
    "Latencia en segundos de /query/llm",
)

# Flujos por tipo (R/C)
query_flow_counter = Counter(
    "query_flow_requests_total",
    "Total requests per query flow",
    ["flow", "status"],
)
query_flow_latency = Histogram(
    "query_flow_request_latency_seconds",
    "Latency per query flow",
    ["flow"],
)

# Noticias / web
news_llm_query_counter = Counter(
    "news_llm_query_requests_total",
    "Total News LLM query requests",
)

news_llm_query_requests_by_status_total = Counter(
    "news_llm_query_requests_by_status_total",
    "News LLM query requests by status",
    ["status"],
)

news_llm_query_latency = Histogram(
    "news_llm_query_request_latency_seconds",
    "Latency of News LLM query requests",
)


#REDIS
def _conv_current_key(user_id: int) -> str:
    return f"conversation_current:{user_id}"


def _conv_legacy_key(user_id: int) -> str:
    # Legacy: JSON con conversation_history mezclado por user
    return f"conversation:{user_id}"


def _conv_meta_key(user_id: int, conversation_id: int) -> str:
    return f"conversation_meta:{user_id}:{conversation_id}"


def _conv_hist_key(user_id: int, conversation_id: int) -> str:
    return f"conversation_history:{user_id}:{conversation_id}"


def _eph_key(user_id: int, conversation_id: int) -> str:
    return f"ephemeral:{user_id}:{conversation_id}"


def _ctx_key(user_id: int, conversation_id: int) -> str:
    return f"context:{user_id}:{conversation_id}"


def _safe_json_loads(raw: Any, default: Any) -> Any:
    if raw is None:
        return default
    try:
        if isinstance(raw, (bytes, bytearray)):
            raw = raw.decode("utf-8", errors="ignore")
        return json.loads(raw)
    except Exception:
        return default


async def _redis_pipeline_execute(pipe) -> None:
    """
    Compatibilidad con redis.asyncio / aioredis.
    """
    res = pipe.execute()
    if asyncio.iscoroutine(res):
        await res


def _now_iso() -> str:
    return datetime.utcnow().isoformat()

#  EXTRACCIÓN DE TEXTO DE FICHEROS
def extract_text_from_pdf(file_path: str) -> str:
    return pdf_extract_text(file_path)


def extract_text_from_docx(file_path: str) -> str:
    return docx2txt.process(file_path)


def extract_text_from_doc(file_path: str) -> str:
    # Si en tu stack usas antiword u otra lib específica, cámbialo aquí.
    return docx2txt.process(file_path)


def extract_text_from_xlsx(file_path: str) -> str:
    df = pd.read_excel(file_path)
    records = df.astype(str).apply(lambda x: " ".join(x), axis=1).tolist()
    return "\n".join(records)


def extract_text_from_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    return "\n".join(df.astype(str).values.flatten())


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def allowed_file(filename: str) -> bool:
    """
    Comprueba si la extensión del fichero está soportada para subir
    como 'ephemeral file'.
    """
    return "." in filename and filename.rsplit(".", 1)[1].lower() in {
        "pdf",
        "pptx",
        "xlsx",
        "docx",
        "doc",
        "csv",
        "jpg",
        "jpeg",
        "png",
        "txt",
    }


# =============================================================================
#  LIMPIEZA DE RESPUESTAS DEL LLM
# =============================================================================

def clean_text(text: str) -> str:
    """
    Limpieza agresiva de la salida del LLM antes de devolverla al usuario:

    - Recorta espacios iniciales/finales.
    - Elimina fences ```...``` y etiquetas ```json.
    - Corta cualquier sección de razonamiento interno tipo "Thought:", "Razonamiento:", etc.
    - Quita disclaimers típicos.
    - Colapsa saltos de línea excesivos.
    """
    if not text:
        return ""

    t = str(text).strip()

    # 1) Quitar fences ```lang ... ```
    if t.startswith("```") and t.endswith("```"):
        t = re.sub(r"^```[a-zA-Z0-9_\-]*\s*\n?", "", t)
        t = re.sub(r"\n?```$", "", t)

    # 2) Cortar a partir de "Thought:" / "Razonamiento:" / "Explanation:"
    thought_patterns = [
        r"\n\s*Thought\s*:.*",
        r"\n\s*Razonamiento\s*:.*",
        r"\n\s*Explanation\s*:.*",
    ]
    for pat in thought_patterns:
        t = re.sub(pat, "", t, flags=re.IGNORECASE | re.DOTALL)

    # 3) Quitar cabeceras tipo "Respuesta:" al inicio
    t = re.sub(
        r"^(respuesta:|respuesta\s*|assistant:|asistente:)\s*",
        "",
        t,
        flags=re.IGNORECASE,
    )

    # 4) Eliminar disclaimers típicos
    disclaimer_patterns = [
        r"como modelo de lenguaje[^.\n]*\.",
        r"como ia(?: de)? lenguaje[^.\n]*\.",
    ]
    for pat in disclaimer_patterns:
        t = re.sub(pat, "", t, flags=re.IGNORECASE)

    # 5) Colapsar >2 saltos de línea
    t = re.sub(r"\n{3,}", "\n\n", t)

    return t.strip()


# =============================================================================
#  RAG (llamada al microservicio NLP)
# =============================================================================

async def perform_rag_search(
    prompt: str,
    access_token: str,
    flow: str = "C",
) -> Dict[str, Any]:
    """
    Llama al microservicio NLP (/search) para obtener resultados RAG.

    Ahora devuelve SIEMPRE un diccionario enriquecido con:
      - status: "ok" | "no_index" | "no_results" | ...
      - results: lista de resultados RAG (cada uno con text, doc_name, meta, similar_blocks, ...)
      - aggregation: resultado de aggregate_excel (conteos), o None
      - used_department: directorio departamental usado, o None
      - cleaned_query: query normalizada usada en el buscador
      - message: mensaje opcional del backend

    Esto permite que el modelo de negocio y el orquestador aprovechen
    no solo los fragmentos, sino también:
      * conteos tabulares (aggregate_excel)
      * información de en qué índice se buscó
      * query limpia usada por el NLP
    """
    flow_norm = (flow or "C").upper()

    payload: Dict[str, Any] = {"query": prompt}
    if flow_norm == "R":
        # Modo rápido → por defecto un top_k más agresivo
        payload["top_k"] = 3

    async with httpx.AsyncClient(timeout=180.0) as client:
        resp = await client.post(
            f"{NLP_APP_URL.rstrip('/')}/search",
            json=payload,
            cookies={"access_token": f"Bearer {access_token}"},
        )

    resp.raise_for_status()
    data = resp.json() or {}

    status = data.get("status") or "ok"

    # Log de control, especialmente útil para no_index / no_results
    logging.info(
        "perform_rag_search: status=%s, msg=%s, used_department=%s, cleaned_query=%s",
        status,
        data.get("message"),
        data.get("used_department"),
        data.get("cleaned_query"),
    )

    # Normalizamos la estructura que devolvemos SIEMPRE
    results = data.get("results") or []
    return {
        "status": status,
        "results": results,
        "aggregation": data.get("aggregation"),
        "used_department": data.get("used_department"),
        "cleaned_query": data.get("cleaned_query") or prompt,
        "message": data.get("message"),
    }



def perform_rag_search_sync(
    prompt: str,
    access_token: str,
    flow: str = "C",
) -> List[dict]:
    """
    Versión síncrona de perform_rag_search, pensada para ser usada desde tools de CrewAI.

    Misma semántica que la versión async.
    """
    if flow == "R":
        payload: Dict[str, Any] = {
            "query": prompt,
            "top_k": 3,
            "retrieve_context": True,
            "top_k_context": 1,
        }
    else:
        payload = {"query": prompt}

    try:
        with httpx.Client(timeout=180.0) as client:
            resp = client.post(
                f"{NLP_APP_URL.rstrip('/')}/search",
                json=payload,
                cookies={"access_token": f"Bearer {access_token}"},
            )
        resp.raise_for_status()
    except httpx.HTTPError as exc:
        logger.error(
            "Error HTTP llamando al servicio NLP RAG: %s (tipo=%s)",
            exc,
            type(exc).__name__,
        )
        raise

    try:
        data = resp.json()
    except ValueError:
        logger.error(
            "Respuesta de NLP RAG no es JSON. status=%s body=%s",
            resp.status_code,
            resp.text[:2000],
        )
        raise

    return data.get("results", []) or []


# =============================================================================
#  JWT / AUTH BÁSICO
# =============================================================================

def verify_token(request: Request) -> Dict[str, Any]:
    """
    Verifica el JWT presente en la cookie 'access_token' o en el header Authorization.
    Lanza HTTPException(403) si falta o es inválido/expirado.
    Retorna el payload en caso de éxito.
    """
    JWT_VERIFY_REQUESTS.inc()
    start_time = time.time()

    access_token = request.cookies.get("access_token")
    token_str: Optional[str] = None

    if access_token:
        logger.info("✱ Extracting access token from session cookie...")
        token_str = access_token.replace("Bearer ", "").strip()
    else:
        auth = request.headers.get("Authorization")
        if auth and auth.lower().startswith("bearer "):
            logger.info("✱ Extracting access token from Authorization header...")
            token_str = auth[7:].strip()

    if not token_str:
        JWT_MISSING.inc()
        JWT_VERIFY_DURATION.observe(time.time() - start_time)
        logger.debug("✱ No token present in cookie/header.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="No autorizado: Token no presente.",
        )

    try:
        logger.info("✱ Decoding token to validate access...")
        payload = jwt.decode(token_str, SECRET_KEY, algorithms=[ALGORITHM])

        email = payload.get("sub")
        user_id = payload.get("user_id")
        if not email or user_id is None:
            JWT_INVALID.inc()
            JWT_VERIFY_DURATION.observe(time.time() - start_time)
            logger.debug("✱ Token payload missing 'sub' or 'user_id'.")
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Token no válido.",
            )

        JWT_VALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - start_time)
        logger.debug(
            "✱ Token verified successfully for user '%s' (id=%s).", email, user_id
        )
        return payload

    except JWTError:
        JWT_INVALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - start_time)
        logger.debug("✱ Token inválido o expirado.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Token no válido o expirado.",
        )


def verify_token_from_cookie(request: Request) -> str:
    """
    Verifica el JWT en la cookie 'access_token'.
    Retorna el email (campo 'sub') si es válido.
    Lanza HTTPException(403) en caso contrario.
    """
    JWT_VERIFY_REQUESTS.inc()
    t0 = time.time()

    access_token = request.cookies.get("access_token")
    if access_token is None:
        JWT_MISSING.inc()
        JWT_VERIFY_DURATION.observe(time.time() - t0)
        logger.debug("✱ Token no presente en la cookie.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="No autorizado: Token no presente.",
        )

    token_str = access_token.replace("Bearer ", "")
    try:
        payload = jwt.decode(token_str, SECRET_KEY, algorithms=[ALGORITHM])
        email = payload.get("sub")
        if not email:
            JWT_INVALID.inc()
            JWT_VERIFY_DURATION.observe(time.time() - t0)
            logger.debug("✱ Token no contiene 'sub'.")
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Token no válido.",
            )

        JWT_VALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - t0)
        return email

    except JWTError:
        JWT_INVALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - t0)
        logger.debug("✱ Token inválido o expirado.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Token no válido o expirado.",
        )


# =============================================================================
#  REDIS CLIENTS (inyectados desde main.py)
# =============================================================================

redis_client = None
redis_client_conversations = None


def init_redis_clients(core_client, conv_client) -> None:
    """
    Registra los clientes Redis que usará el microservicio de modelo de negocio.

    Debe llamarse desde main.py en el evento startup.
    """
    global redis_client, redis_client_conversations
    redis_client = core_client
    redis_client_conversations = conv_client
    logger.info(
        "init_redis_clients: Redis clients registrados (core db=0, conv db posiblemente 2)."
    )


# =============================================================================
#  CSRF (double-submit cookie)
# =============================================================================

def generate_csrf_token() -> str:
    """
    Genera un token CSRF aleatorio y seguro.
    """
    return secrets.token_urlsafe(32)


def validate_csrf_double_submit(
    request: Request,
    cookie_name: str,
    header_name: str = "X-CSRFToken",
    error_detail: str = "CSRF token inválido o ausente.",
) -> None:
    """
    Valida el token CSRF usando el patrón de double-submit cookie:
      - Cookie: cookie_name
      - Cabecera: header_name
    Deben coincidir.
    """
    cookie_token = request.cookies.get(cookie_name)
    header_token = request.headers.get(header_name)

    if not cookie_token or not header_token or cookie_token != header_token:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail=error_detail,
        )


# =============================================================================
#  HELPERS DE AUTH / COSTE
# =============================================================================

def extract_raw_bearer_token(request: Request) -> Optional[str]:
    """
    Extrae el JWT crudo (sin el prefijo 'Bearer ') desde:

      1) Cabecera Authorization: 'Bearer <jwt>'
      2) Cookie 'access_token': 'Bearer <jwt>'
    """
    auth_header = request.headers.get("Authorization", "")
    cookie_token = request.cookies.get("access_token")

    if auth_header.startswith("Bearer "):
        return auth_header.split(" ", 1)[1].strip()

    if cookie_token and cookie_token.startswith("Bearer "):
        return cookie_token.split(" ", 1)[1].strip()

    return None


def calculate_cost(total_tokens: int, price_per_1k: float = 0.0005) -> float:
    """
    Coste aproximado: price_per_1k por cada 1000 tokens.
    Se guarda como float en UsageLog.cost (Numeric(10,6)).
    """
    return round((total_tokens / 1000.0) * price_per_1k, 6)


# =============================================================================
#  CONVERSACIONES (BD)
# =============================================================================

def get_or_create_conversation(
    db: Session,
    user_id: int,
    conversation_id: Optional[int],
) -> Tuple[Conversation, bool]:
    """
    Devuelve una conversación existente del usuario o crea una nueva.

    return: (conversation, created)
      - created=True si se ha creado una nueva conversación.
      - created=False si se ha reutilizado una existente.
    """
    if conversation_id:
        convo = (
            db.query(Conversation)
            .filter(Conversation.id == conversation_id, Conversation.user_id == user_id)
            .first()
        )
        if convo:
            return convo, False

    convo = Conversation(
        user_id=user_id,
        conversation_text="",
        created_at=datetime.utcnow(),
    )
    db.add(convo)
    db.commit()
    db.refresh(convo)
    return convo, True


# =============================================================================
#  REDIS: SESIÓN + CONVERSACIÓN + ARCHIVOS EFÍMEROS
# =============================================================================

async def get_session_from_redis(user_id: int) -> Optional[Dict[str, Any]]:
    """
    Recupera la sesión del usuario desde Redis (clave session:{user_id}).
    """
    logger.info(
        "\n\n(utils - get_session_from_redis)  Llamada asíncrona a REDIS para obtener los datos de sesión..."
    )
    REDIS_SESSION_REQUESTS.inc()
    start = time.time()

    if redis_client is None:
        logger.error("get_session_from_redis: redis_client no inicializado.")
        REDIS_SESSION_MISSES.inc()
        REDIS_SESSION_DURATION.observe(time.time() - start)
        return None

    session = await redis_client.get(f"session:{user_id}")
    duration = time.time() - start
    REDIS_SESSION_DURATION.observe(duration)

    if session:
        REDIS_SESSION_HITS.inc()
        try:
            return json.loads(session)
        except json.JSONDecodeError:
            logger.warning(
                "get_session_from_redis: JSON corrupto para session:%s", user_id
            )
            REDIS_SESSION_MISSES.inc()
            return None

    REDIS_SESSION_MISSES.inc()
    return None


async def save_conversation_to_redis(
    user_id: int,
    expires_in: int = 480,
    conversation_entry: Optional[dict] = None,
    conversation_id: Optional[int] = None,
) -> None:
    """
    Guarda el historial de conversación en Redis aislado por conversación.

    NUEVO (sin romper):
      - Si se pasa conversation_id (recomendado): usa claves por conversación.
      - Si NO se pasa: intenta inferirlo desde conversation_entry['conversation_id']
        o desde conversation_current:{user_id}. Si no puede, cae a legacy.
    """
    import uuid  # <- FIX: evita NameError si el módulo no importó uuid

    logger.info(
        "\n\n(utils - save_conversation_to_redis)  Actualizando el historial de conversación en Redis para el usuario..."
    )
    CONV_SAVE_REQUESTS.inc()
    start = time.time()

    if redis_client_conversations is None:
        logger.error("save_conversation_to_redis: redis_client_conversations no inicializado.")
        CONV_SAVE_DURATION.observe(time.time() - start)
        return

    # 1) Determinar conversation_id scope
    cid = conversation_id
    if cid is None and isinstance(conversation_entry, dict):
        ce_cid = conversation_entry.get("conversation_id")
        if ce_cid is not None:
            try:
                cid = int(ce_cid)
            except Exception:
                cid = None

    if cid is None:
        # Best-effort: puntero de "conversación actual" por user
        current = await redis_client_conversations.get(_conv_current_key(user_id))
        if current:
            try:
                if isinstance(current, (bytes, bytearray)):
                    current = current.decode("utf-8", errors="ignore")
                cid = int(str(current).strip())
            except Exception:
                cid = None

    # 2) Si no tenemos cid, caemos al legacy (para no romper llamadas antiguas)
    if cid is None:
        logger.warning(
            "save_conversation_to_redis: sin conversation_id; usando clave legacy conversation:%s (puede mezclar contextos).",
            user_id,
        )

        legacy_key = _conv_legacy_key(user_id)
        conversation = await redis_client_conversations.get(legacy_key)
        conversation_data = _safe_json_loads(conversation, default=None)

        if not isinstance(conversation_data, dict):
            conversation_data = {
                "conversation_id": str(uuid.uuid4()),
                "user_id": user_id,
                "conversation_history": [],
            }

        if "conversation_history" not in conversation_data or not isinstance(conversation_data["conversation_history"], list):
            conversation_data["conversation_history"] = []

        if conversation_entry and isinstance(conversation_entry, dict) and "conversation_id" in conversation_entry:
            conversation_data["current_conversation_id"] = conversation_entry["conversation_id"]

        conversation_data["conversation_history"].append(conversation_entry or {})

        await redis_client_conversations.set(
            legacy_key,
            json.dumps(conversation_data, ensure_ascii=False),
            ex=expires_in,
        )

        CONV_SAVE_DURATION.observe(time.time() - start)
        return

    # 3) Camino nuevo: por conversación
    meta_key = _conv_meta_key(user_id, cid)
    hist_key = _conv_hist_key(user_id, cid)
    current_key = _conv_current_key(user_id)

    entry_obj = conversation_entry or {}
    try:
        entry_json = json.dumps(entry_obj, ensure_ascii=False)
    except Exception:
        entry_json = json.dumps({"_raw": str(entry_obj)}, ensure_ascii=False)

    max_entries = int(os.getenv("CONVERSATION_HISTORY_MAX_ENTRIES", "60"))

    pipe = redis_client_conversations.pipeline(transaction=True)

    # A) set puntero actual por user (best-effort)
    pipe.set(current_key, str(cid), ex=expires_in)

    # B) crear/actualizar meta
    existing_meta_raw = await redis_client_conversations.get(meta_key)
    existing_meta = _safe_json_loads(existing_meta_raw, default=None)

    if not isinstance(existing_meta, dict):
        meta_obj = {
            "conversation_uuid": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_id": cid,
            "created_at": _now_iso(),
            "updated_at": _now_iso(),
        }
        pipe.set(meta_key, json.dumps(meta_obj, ensure_ascii=False), ex=expires_in)
    else:
        existing_meta["updated_at"] = _now_iso()
        pipe.set(meta_key, json.dumps(existing_meta, ensure_ascii=False), ex=expires_in)

    # C) append a LIST (atómico en Redis)
    pipe.rpush(hist_key, entry_json)

    # D) trim para evitar crecimiento infinito (mantener últimos N)
    if max_entries > 0:
        pipe.ltrim(hist_key, max(-max_entries, -1000000), -1)

    # E) TTL
    pipe.expire(hist_key, expires_in)

    await _redis_pipeline_execute(pipe)

    CONV_SAVE_DURATION.observe(time.time() - start)


async def add_ephemeral_file(
    user_id: int,
    filename: str,
    text: str,
    ttl: int = 3600,
    meta: Optional[Dict[str, Any]] = None,
    conversation_id: Optional[int] = None,
) -> None:
    """
    Añade/actualiza archivos en vuelo en Redis.

    NUEVO (sin romper):
      - Si se pasa conversation_id: ephemeral:{user_id}:{conversation_id}
      - Si NO se pasa: intenta inferir conversation_id desde conversation_current:{user_id}
        y usar clave scoped.
      - Si no puede inferir, cae a legacy ephemeral:{user_id}.
    """
    if redis_client is None:
        logger.error("add_ephemeral_file: redis_client no inicializado.")
        return

    cid = None
    if conversation_id is not None:
        try:
            cid = int(conversation_id)
        except Exception:
            cid = None

    # Best-effort: si no llega cid, intentamos usar puntero actual
    if cid is None and redis_client_conversations is not None:
        current = await redis_client_conversations.get(_conv_current_key(user_id))
        if current:
            try:
                if isinstance(current, (bytes, bytearray)):
                    current = current.decode("utf-8", errors="ignore")
                cid = int(str(current).strip())
            except Exception:
                cid = None

    if cid is None:
        key = f"ephemeral:{user_id}"  # legacy (puede mezclar)
    else:
        key = _eph_key(user_id, cid)

    existing = await redis_client.get(key)
    files = _safe_json_loads(existing, default=[])
    if not isinstance(files, list):
        files = []

    max_stored_chars = int(os.getenv("EPHEMERAL_MAX_STORED_CHARS", "12000"))
    max_files = int(os.getenv("EPHEMERAL_MAX_FILES", "25"))

    record = {
        "filename": filename,
        "text": (text or "")[:max_stored_chars],
        "uploaded_at": datetime.utcnow().isoformat(),
        "meta": meta or {},
        "conversation_id": cid if cid is not None else conversation_id,
    }

    files.append(record)

    # Limitar número total de archivos guardados (mantener los últimos N)
    if max_files > 0 and len(files) > max_files:
        files = files[-max_files:]

    await redis_client.set(key, json.dumps(files, ensure_ascii=False), ex=ttl)


async def get_ephemeral_files(
    user_id: int,
    conversation_id: Optional[int] = None,
    file_ids: Optional[List[str]] = None,
) -> List[dict]:
    """
    Devuelve la lista de archivos en vuelo del usuario.

    NUEVO (sin romper):
      - Si se pasa conversation_id: lee ephemeral:{user_id}:{conversation_id}
      - Si NO se pasa: intenta usar conversation_current:{user_id} para leer scoped
        y, si no puede, cae a legacy ephemeral:{user_id}.
      - Si file_ids se pasa, filtra por meta.file_id (o file_id top-level).
    """
    if redis_client is None:
        logger.error("get_ephemeral_files: redis_client no inicializado.")
        return []

    cid = None
    if conversation_id is not None:
        try:
            cid = int(conversation_id)
        except Exception:
            cid = None

    if cid is None and redis_client_conversations is not None:
        current = await redis_client_conversations.get(_conv_current_key(user_id))
        if current:
            try:
                if isinstance(current, (bytes, bytearray)):
                    current = current.decode("utf-8", errors="ignore")
                cid = int(str(current).strip())
            except Exception:
                cid = None

    if cid is None:
        key = f"ephemeral:{user_id}"  # legacy
    else:
        key = _eph_key(user_id, cid)

    existing = await redis_client.get(key)
    if not existing:
        return []

    files = _safe_json_loads(existing, default=[])
    if not isinstance(files, list):
        logger.warning("get_ephemeral_files: JSON corrupto o tipo inválido para key=%s", key)
        return []

    wanted = set(file_ids or [])
    out: List[dict] = []

    for item in files:
        if not isinstance(item, dict):
            continue

        # Seguridad: si el item viene con conversation_id y NO coincide con cid, lo ignoramos
        item_cid = item.get("conversation_id")
        if cid is not None and item_cid is not None:
            try:
                if int(item_cid) != int(cid):
                    continue
            except Exception:
                # si es inválido, mejor ignorar el item
                continue

        if wanted:
            meta = item.get("meta") or {}
            fid = meta.get("file_id") or item.get("file_id")
            if fid not in wanted:
                continue

        out.append(item)

    return out


async def get_conversation_to_redis(
    user_id: int,
    expires_in: int = 480,
    conversation_id: Optional[int] = None,
) -> dict:
    """
    Recupera la conversación en Redis.

    NUEVO (sin romper):
      - Si se pasa conversation_id: devuelve el historial de ESA conversación (aislado).
      - Si no: intenta usar conversation_current:{user_id} o legacy.
    """
    logger.info(
        "\n\n(utils - get_conversation_to_redis)  Recuperando la conversación en Redis para el usuario..."
    )
    CONV_GET_REQUESTS.inc()
    start = time.time()

    if redis_client_conversations is None:
        logger.error("get_conversation_to_redis: redis_client_conversations no inicializado.")
        CONV_GET_MISSES.inc()
        CONV_GET_DURATION.observe(time.time() - start)
        return {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "current_conversation_id": conversation_id,
            "conversation_history": [],
        }

    cid = conversation_id
    if cid is None:
        current = await redis_client_conversations.get(_conv_current_key(user_id))
        if current:
            try:
                if isinstance(current, (bytes, bytearray)):
                    current = current.decode("utf-8", errors="ignore")
                cid = int(str(current).strip())
            except Exception:
                cid = None

    # Si aún no hay cid, caer a legacy (para no romper)
    if cid is None:
        legacy_key = _conv_legacy_key(user_id)
        conversation = await redis_client_conversations.get(legacy_key)
        if conversation:
            CONV_GET_HITS.inc()
            conversation_data = _safe_json_loads(conversation, default=None)
            if not isinstance(conversation_data, dict):
                conversation_data = {
                    "conversation_id": str(uuid.uuid4()),
                    "user_id": user_id,
                    "conversation_history": [],
                }
        else:
            CONV_GET_MISSES.inc()
            conversation_data = {
                "conversation_id": str(uuid.uuid4()),
                "user_id": user_id,
                "conversation_history": [],
            }
            await redis_client_conversations.set(
                legacy_key,
                json.dumps(conversation_data, ensure_ascii=False),
                ex=expires_in,
            )

        CONV_GET_DURATION.observe(time.time() - start)
        return conversation_data

    meta_key = _conv_meta_key(user_id, cid)
    hist_key = _conv_hist_key(user_id, cid)

    meta_raw = await redis_client_conversations.get(meta_key)
    meta = _safe_json_loads(meta_raw, default=None)

    # ---- Migración best-effort desde legacy, filtrando por conversation_id si hay entries
    if not isinstance(meta, dict):
        legacy_key = _conv_legacy_key(user_id)
        legacy_raw = await redis_client_conversations.get(legacy_key)
        legacy = _safe_json_loads(legacy_raw, default=None)

        # Creamos meta nueva
        meta = {
            "conversation_uuid": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_id": cid,
            "created_at": _now_iso(),
            "updated_at": _now_iso(),
            "migrated_from_legacy": False,
        }

        # Si legacy trae history, intentamos migrar SOLO entradas de esta conversación (si viene marcado)
        if isinstance(legacy, dict):
            legacy_hist = legacy.get("conversation_history")
            if isinstance(legacy_hist, list) and legacy_hist:
                filtered = []
                for item in legacy_hist:
                    if not isinstance(item, dict):
                        continue
                    item_cid = item.get("conversation_id")
                    try:
                        if item_cid is not None and int(item_cid) == int(cid):
                            filtered.append(item)
                    except Exception:
                        continue

                if filtered:
                    # Pasa filtered a LIST si aún no hay historial en LIST
                    existing_len = await redis_client_conversations.llen(hist_key)
                    if existing_len == 0:
                        pipe = redis_client_conversations.pipeline(transaction=True)
                        for it in filtered:
                            pipe.rpush(hist_key, json.dumps(it, ensure_ascii=False))
                        max_entries = int(os.getenv("CONVERSATION_HISTORY_MAX_ENTRIES", "60"))
                        if max_entries > 0:
                            pipe.ltrim(hist_key, max(-max_entries, -1000000), -1)
                        pipe.expire(hist_key, expires_in)
                        await _redis_pipeline_execute(pipe)
                        meta["migrated_from_legacy"] = True

        await redis_client_conversations.set(
            meta_key, json.dumps(meta, ensure_ascii=False), ex=expires_in
        )

    # Leemos historial (LIST) y devolvemos el mismo shape que antes
    max_entries = int(os.getenv("CONVERSATION_HISTORY_MAX_ENTRIES", "60"))
    if max_entries > 0:
        raw_items = await redis_client_conversations.lrange(hist_key, -max_entries, -1)
    else:
        raw_items = await redis_client_conversations.lrange(hist_key, 0, -1)

    history: List[dict] = []
    for raw in raw_items or []:
        obj = _safe_json_loads(raw, default=None)
        if isinstance(obj, dict):
            history.append(obj)

    # Refrescamos TTLs best-effort
    pipe = redis_client_conversations.pipeline(transaction=True)
    pipe.expire(hist_key, expires_in)
    pipe.expire(meta_key, expires_in)
    pipe.set(_conv_current_key(user_id), str(cid), ex=expires_in)
    await _redis_pipeline_execute(pipe)

    CONV_GET_HITS.inc()
    CONV_GET_DURATION.observe(time.time() - start)

    return {
        "conversation_id": meta.get("conversation_uuid") or str(uuid.uuid4()),
        "user_id": user_id,
        "current_conversation_id": cid,
        "conversation_history": history,
        "meta": meta,
    }

async def reset_conversation_in_redis(
    user_id: int,
    expires_in: int = 480,
    conversation_id: Optional[int] = None,
) -> dict:
    """
    Resetea la conversación asociada al usuario en Redis.

    NUEVO (sin romper):
      - Si se pasa conversation_id: resetea SOLO esa conversación (historial/lista + meta)
        y limpia también eph/ctx scoped (en redis_client).
      - Si no: resetea legacy conversation:{user_id} (comportamiento anterior)
        y limpia legacy eph/ctx best-effort.
    """
    import uuid  # por seguridad

    logger.info(
        "\n\n(utils - reset_conversation_in_redis)  Reseteando la conversación en Redis para el usuario..."
    )
    CONV_RESET_REQUESTS.inc()
    start = time.time()

    if redis_client_conversations is None:
        logger.error("reset_conversation_in_redis: redis_client_conversations no inicializado.")
        CONV_RESET_DURATION.observe(time.time() - start)
        return {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_history": [],
        }

    if conversation_id is None:
        # Legacy reset (mezcla por user)
        new_conversation_data = {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_history": [],
        }
        await redis_client_conversations.set(
            _conv_legacy_key(user_id),
            json.dumps(new_conversation_data, ensure_ascii=False),
            ex=expires_in,
        )
        await redis_client_conversations.delete(_conv_current_key(user_id))

        # Best-effort: limpiar eph/ctx legacy
        if redis_client is not None:
            try:
                await redis_client.delete(f"ephemeral:{user_id}", f"context:{user_id}")
            except Exception:
                pass

        CONV_RESET_DURATION.observe(time.time() - start)
        return new_conversation_data

    cid = int(conversation_id)
    meta_key = _conv_meta_key(user_id, cid)
    hist_key = _conv_hist_key(user_id, cid)

    meta_obj = {
        "conversation_uuid": str(uuid.uuid4()),
        "user_id": user_id,
        "conversation_id": cid,
        "created_at": _now_iso(),
        "updated_at": _now_iso(),
        "reset": True,
    }

    pipe = redis_client_conversations.pipeline(transaction=True)
    pipe.delete(hist_key)
    pipe.set(meta_key, json.dumps(meta_obj, ensure_ascii=False), ex=expires_in)
    pipe.set(_conv_current_key(user_id), str(cid), ex=expires_in)
    await _redis_pipeline_execute(pipe)

    # FIX CLAVE: limpiar efímeros/contexto scoped en el Redis "core"
    if redis_client is not None:
        try:
            await redis_client.delete(_eph_key(user_id, cid), _ctx_key(user_id, cid))
        except Exception:
            # no hacemos fail hard: reset conversacional debe seguir funcionando
            pass

    CONV_RESET_DURATION.observe(time.time() - start)
    return {
        "conversation_id": meta_obj["conversation_uuid"],
        "user_id": user_id,
        "current_conversation_id": cid,
        "conversation_history": [],
        "meta": meta_obj,
    }


async def set_user_context(
    user_id: int,
    text: str,
    ttl: int = 3600,
    conversation_id: Optional[int] = None,
) -> None:
    """
    Guarda el último contexto textual.

    NUEVO (sin romper):
      - Si conversation_id: context:{user_id}:{conversation_id}
      - Si NO: intenta inferir conversation_current:{user_id} y guardar scoped.
      - Si no puede, cae a legacy context:{user_id}.
    """
    if redis_client is None:
        logger.error("set_user_context: redis_client no inicializado.")
        return

    cid = None
    if conversation_id is not None:
        try:
            cid = int(conversation_id)
        except Exception:
            cid = None

    if cid is None and redis_client_conversations is not None:
        current = await redis_client_conversations.get(_conv_current_key(user_id))
        if current:
            try:
                if isinstance(current, (bytes, bytearray)):
                    current = current.decode("utf-8", errors="ignore")
                cid = int(str(current).strip())
            except Exception:
                cid = None

    if cid is None:
        key = f"context:{user_id}"  # legacy
    else:
        key = _ctx_key(user_id, cid)

    await redis_client.set(key, text, ex=ttl)


# =============================================================================
#  AUTH COMPUESTO: SESIÓN UNÍVOCA (JWT + REDIS)
# =============================================================================

async def get_current_auth(request: Request) -> Dict[str, Any]:
    """
    Dependencia central de autenticación para el microservicio de modelo de negocio.

    Flujo:
      - verify_token valida el JWT y devuelve el payload (token).
      - extract_raw_bearer_token saca el JWT crudo de la petición.
      - get_session_from_redis recupera la sesión session:{user_id}.
      - Se compara session['token'] con el JWT crudo -> sesión unívoca.
    """
    # 1) Validar JWT (cookie o header)
    token = verify_token(request)
    user_id = token.get("user_id")
    if not user_id:
        logger.warning("get_current_auth: token sin user_id.")
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Token inválido (sin user_id).",
        )

    # 2) Extraer JWT "crudo" enviado en la petición
    raw_token = extract_raw_bearer_token(request)
    if not raw_token:
        logger.warning(
            "get_current_auth: no se pudo extraer token Bearer de la petición."
        )
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Credenciales de autenticación no proporcionadas.",
        )

    # 3) Recuperar sesión desde Redis
    session_data = await get_session_from_redis(user_id)
    if not session_data:
        logger.info(
            "get_current_auth: sesión no encontrada en Redis para user_id=%s.",
            user_id,
        )
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Sesión expirada o no encontrada.",
        )

    # 4) Comprobar que el token coincide con el de la sesión
    session_token = session_data.get("token")
    if session_token != raw_token:
        logger.info(
            "get_current_auth: token de la petición no coincide con la sesión activa "
            "en Redis para user_id=%s. Posible sesión revocada o reemplazada.",
            user_id,
        )
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Sesión revocada o reemplazada. Por favor, inicia sesión de nuevo.",
        )

    return {
        "user_id": user_id,
        "role": token.get("role"),
        "token_payload": token,
        "access_token": raw_token,
        "session": session_data,
    }


# -------------------------------------------------
# CLAMAV
# -------------------------------------------------
async def scan_with_clamav(content: bytes, filename: Optional[str] = None) -> Dict[str, Any]:
    """
    Escanea un blob de bytes con ClamAV (clamd) vía INSTREAM.

    Devuelve un diccionario con:
      {
        "status": "OK" | "INFECTED" | "ERROR" | "SKIPPED",
        "virus_name": str | None,
        "raw_result": str | None,
        "error": str | None,
        "duration_s": float,
        "filename": str | None,
        "size_bytes": int
      }

    Política por defecto:
      - Si ClamAV responde 'FOUND' → status='INFECTED'
      - Si ClamAV responde 'OK'    → status='OK'
      - Cualquier otro resultado o excepción → status='ERROR'
    """
    start = time.time()

    # Fichero vacío → lo damos por OK
    if not content:
        return {
            "status": "OK",
            "virus_name": None,
            "raw_result": "EMPTY_FILE",
            "error": None,
            "duration_s": time.time() - start,
            "filename": filename,
            "size_bytes": 0,
        }

    # Host / puerto de ClamAV:
    # - En docker-compose, normalmente CLAMAV_HOST=clamav, CLAMAV_PORT=3310
    # - En host Windows con clamd publicado: CLAMAV_HOST=127.0.0.1
    host = os.getenv("CLAMAV_HOST", "127.0.0.1")
    port = int(os.getenv("CLAMAV_PORT", "3310"))
    timeout = float(os.getenv("CLAMAV_TIMEOUT", "60.0"))

    # Prefijo de comando: 'z' (NULL-terminated) o 'n' (newline-terminated).
    # ClamAV suele aceptar 'z' para TCP ('zINSTREAM\\0'), pero se puede ajustar
    # vía variable de entorno CLAMAV_CMD_PREFIX.
    cmd_prefix = os.getenv("CLAMAV_CMD_PREFIX", "z").lower()
    if cmd_prefix not in ("z", "n"):
        cmd_prefix = "z"

    loop = asyncio.get_running_loop()
    result = await loop.run_in_executor(
        None,
        _scan_with_clamav_sync,
        host,
        port,
        timeout,
        cmd_prefix,
        content,
        filename,
        start,
    )
    return result


def _scan_with_clamav_sync(
    host: str,
    port: int,
    timeout: float,
    cmd_prefix: str,
    content: bytes,
    filename: Optional[str],
    start: float,
) -> Dict[str, Any]:
    """
    Implementación síncrona del escaneo INSTREAM.

    - Intenta conectar primero al host configurado (CLAMAV_HOST).
    - Si falla la conexión, prueba hosts típicos en entorno Docker/Windows:
        * clamav
        * clamav_service
        * host.docker.internal
        * 127.0.0.1
        * localhost
    - Envía directamente el comando INSTREAM sin PING previo para evitar
      que algunos servidores cierren la conexión.
    """
    size_bytes = len(content)
    raw_result: Optional[str] = None
    virus_name: Optional[str] = None
    error: Optional[str] = None
    status = "ERROR"

    display_filename = filename or "<sin_nombre>"

    sock: Optional[socket.socket] = None
    last_exc: Optional[BaseException] = None

    # Fallbacks de host para soportar:
    #   - modelo de negocio en contenedor Docker
    #   - ClamAV en otro contenedor o en el host
    env_fallbacks = os.getenv("CLAMAV_HOST_FALLBACKS", "")
    extra_hosts: List[str] = [
        h.strip()
        for h in env_fallbacks.split(",")
        if h.strip()
    ]

    candidate_hosts: List[str] = [host] + extra_hosts
    for h in ("clamav", "clamav_service", "host.docker.internal", "127.0.0.1", "localhost"):
        if h not in candidate_hosts:
            candidate_hosts.append(h)

    host_used = host

    try:
        # 1) Conexión a alguno de los hosts candidatos
        for candidate in candidate_hosts:
            try:
                logger.info(
                    "[AV] Conectando a ClamAV en %s:%s para escanear '%s' (%d bytes)...",
                    candidate,
                    port,
                    display_filename,
                    size_bytes,
                )
                sock = socket.create_connection((candidate, port), timeout=timeout)
                host_used = candidate
                break
            except OSError as exc:
                last_exc = exc
                logger.warning(
                    "[AV] No se pudo conectar a ClamAV en %s:%s: %s",
                    candidate,
                    port,
                    exc,
                )

        if sock is None:
            # No se pudo conectar a ningún host candidato
            raise last_exc or RuntimeError(
                f"No se pudo conectar a ClamAV en {host}:{port}"
            )

        # 2) Preparar comando INSTREAM según prefijo z/n
        if cmd_prefix == "z":
            instream_cmd = b"zINSTREAM\0"
        else:  # "n"
            instream_cmd = b"nINSTREAM\n"

        chunk_size = int(os.getenv("CLAMAV_CHUNK_SIZE", "8192"))

        with sock:
            # 2) Protocolo INSTREAM
            # Prefijo INSTREAM + stream de chunks + chunk final de tamaño 0
            sock.sendall(instream_cmd)

            view = memoryview(content)
            offset = 0

            while offset < size_bytes:
                chunk = view[offset : offset + chunk_size]
                # Primero longitud (4 bytes BE) y luego los datos
                sock.sendall(struct.pack(">L", len(chunk)))
                sock.sendall(chunk)
                offset += len(chunk)

            # Fin de stream
            sock.sendall(struct.pack(">L", 0))

            # 3) Leer respuesta (ej: "stream: OK" o "...FOUND")
            chunks: List[bytes] = []
            while True:
                data = sock.recv(4096)
                if not data:
                    break
                chunks.append(data)

        raw_result = b"".join(chunks).decode(errors="replace").strip("\x00\r\n\t ")
        logger.info(
            "[AV] Resultado ClamAV para '%s' (host=%s): %s",
            display_filename,
            host_used,
            raw_result,
        )

        # 4) Interpretar resultado
        if "FOUND" in raw_result:
            status = "INFECTED"
            try:
                # "stream: <nombre_virus> FOUND"
                after_colon = raw_result.split(":", 1)[1].strip()
                virus_name = after_colon.rsplit(" ", 1)[0].strip()
            except Exception:
                virus_name = None
        elif "OK" in raw_result:
            status = "OK"
        else:
            status = "ERROR"
            error = f"Resultado inesperado de ClamAV: {raw_result}"

    except Exception as exc:
        # Cubre:
        #   - ConnectionRefused / timeout
        #   - ConnectionAborted (WinError 10053)
        #   - errores de protocolo
        error = str(exc)

        # Mensaje más explícito para el caso típico en Windows
        if isinstance(exc, ConnectionAbortedError):
            logger.error(
                "[AV] Conexión abortada al escanear '%s' con ClamAV (host=%s:%s): %s. "
                "En Windows esto suele indicar que el servidor cerró la conexión "
                "al recibir un comando inesperado o que un firewall/antivirus local "
                "interrumpió el flujo.",
                display_filename,
                host_used,
                port,
                exc,
                exc_info=True,
            )
        else:
            logger.error(
                "[AV] Error al escanear '%s' con ClamAV (host=%s:%s): %s",
                display_filename,
                host_used,
                port,
                exc,
                exc_info=True,
            )

        status = "ERROR"

    duration = time.time() - start

    return {
        "status": status,
        "virus_name": virus_name,
        "raw_result": raw_result,
        "error": error,
        "duration_s": duration,
        "filename": filename,
        "size_bytes": size_bytes,
    }


def _b64url(data: bytes) -> str:
    return base64.urlsafe_b64encode(data).rstrip(b"=").decode("utf-8")


def _create_hs256_jwt(payload: Dict[str, Any], secret: str) -> str:
    """
    JWT HS256 sin dependencias externas.
    """
    header = {"alg": "HS256", "typ": "JWT"}
    header_b64 = _b64url(json.dumps(header, separators=(",", ":"), ensure_ascii=False).encode("utf-8"))
    payload_b64 = _b64url(json.dumps(payload, separators=(",", ":"), ensure_ascii=False).encode("utf-8"))
    signing_input = f"{header_b64}.{payload_b64}".encode("utf-8")
    signature = hmac.new(secret.encode("utf-8"), signing_input, hashlib.sha256).digest()
    sig_b64 = _b64url(signature)
    return f"{header_b64}.{payload_b64}.{sig_b64}"


# -------------------------
# Config Notetaker URL
# -------------------------
def _get_notetaker_base_url() -> str:
    """
    URL base del frontend de Notetaker. Debe venir por entorno.
    Ej: https://notetaker.cosgs.com
    """
    base = (
        os.getenv("NOTETAKER_FRONTEND_URL")
        or os.getenv("NOTETAKER_URL")
        or os.getenv("REACT_APP_NOTETAKER_URL")
        or os.getenv("VITE_NOTETAKER_URL")
        or ""
    ).strip()

    if not base:
        raise HTTPException(status_code=500, detail="NOTETAKER_FRONTEND_URL no está configurada en el entorno.")

    base = base.rstrip("/")
    parsed = urlparse(base)

    if parsed.scheme not in ("https", "http") or not parsed.netloc:
        raise HTTPException(status_code=500, detail="NOTETAKER_FRONTEND_URL inválida (debe incluir esquema y host).")

    # Hardening: forzar HTTPS salvo entornos controlados
    allow_http = os.getenv("NOTETAKER_ALLOW_HTTP", "false").lower() in ("1", "true", "yes")
    if parsed.scheme == "http" and not allow_http:
        raise HTTPException(
            status_code=500,
            detail="NOTETAKER_FRONTEND_URL debe usar https (o habilita NOTETAKER_ALLOW_HTTP=true para no-prod).",
        )

    return base


# -------------------------
# display_name normalization
# -------------------------
_DISPLAY_NAME_BAD_CHARS = re.compile(r"[\r\n\t]")
_DISPLAY_NAME_HAS_AT = re.compile(r"@")
_DISPLAY_NAME_MULTI_SPACE = re.compile(r"\s+")


def _normalize_display_name(value: Optional[str]) -> str:
    """
    display_name debe ser "Nombre Apellidos" (no email).
    - limpia espacios
    - rechaza si contiene '@'
    - exige al menos 2 palabras
    """
    if not value:
        return ""
    v = " ".join(value.strip().split())
    if not v:
        return ""
    if "@" in v:
        return ""
    if len(v.split()) < 2:
        return ""
    return v


async def _resolve_notetaker_identity_like_me(request: Request, auth: Dict[str, Any]) -> Tuple[str, str]:
    """
    Copia el flujo de /me:
      - intenta SSO /auth/me
      - fallback a session Redis
    Y produce:
      - email (obligatorio)
      - display_name (opcional, solo si parece nombre+apellidos)
        usando el MISMO 'username' que expones en /me.
    """
    session = auth.get("session") or {}

    sso_data: Optional[Dict[str, Any]] = None
    auth_sso_base = os.getenv("AUTH_SSO_INTERNAL_BASE", "").strip() or None

    if auth_sso_base:
        try:
            cookies: Dict[str, str] = {}
            access_cookie = request.cookies.get("access_token")
            if access_cookie:
                cookies["access_token"] = access_cookie

            async with httpx.AsyncClient(timeout=8.0) as client:
                resp = await client.get(f"{auth_sso_base}/auth/me", cookies=cookies)

            if resp.status_code == 200:
                sso_data = resp.json()
            else:
                logger.warning("notetaker: auth_sso /auth/me devolvió código %s", resp.status_code)
        except Exception as e:
            logger.warning("notetaker: error llamando a auth_sso /auth/me: %s", e)

    sso = sso_data or {}

    # email igual que /me
    email = (sso.get("email") or session.get("email", "") or "").strip()
    if not email:
        raise HTTPException(status_code=400, detail="No se pudo determinar el email del usuario.")

    # username igual que /me (ojo: tu /me mete fallback al prefijo del email,
    # aquí NO lo usamos para display_name porque no es nombre+apellidos).
    username = (
        (sso.get("name") or "")  # En tu /me lo usas como username
        or (session.get("username", "") or "")
    ).strip()

    display_name = _normalize_display_name(username)

    return email, display_name


# -------------------------
# Build final Notetaker URL
# -------------------------
def _build_notetaker_url(email: str, display_name: str, user_id: int) -> Dict[str, Any]:
    """
    Construye la URL final a /sso-callback con email y display_name (si existe).
    Opcional: cosmos_token (JWT corto) si está habilitado por entorno.
    """
    base = _get_notetaker_base_url()
    callback = f"{base}/sso-callback"

    include_token = os.getenv("NOTETAKER_INCLUDE_COSMOS_TOKEN", "false").lower() in ("1", "true", "yes")

    ttl_raw = os.getenv("NOTETAKER_COSMOS_TOKEN_TTL_SECONDS", "300")
    try:
        ttl = int(ttl_raw)
    except Exception:
        raise HTTPException(status_code=500, detail="NOTETAKER_COSMOS_TOKEN_TTL_SECONDS debe ser un entero.")
    if ttl < 60 or ttl > 3600:
        # hardening: TTL razonable (1min..1h)
        raise HTTPException(status_code=500, detail="NOTETAKER_COSMOS_TOKEN_TTL_SECONDS fuera de rango (60..3600).")

    params: Dict[str, str] = {"email": email}

    # display_name solo si es válido (ya normalizado)
    if display_name:
        params["display_name"] = display_name

    token_value: Optional[str] = None
    if include_token:
        secret = os.getenv("NOTETAKER_COSMOS_TOKEN_SECRET", "").strip()
        if not secret:
            raise HTTPException(
                status_code=500,
                detail="NOTETAKER_INCLUDE_COSMOS_TOKEN está habilitado pero falta NOTETAKER_COSMOS_TOKEN_SECRET.",
            )

        now = int(datetime.now(timezone.utc).timestamp())
        payload = {
            "iss": "cosmos",
            "sub": str(user_id),
            "email": email,
            # Nota: si display_name está vacío, no lo metemos (evita basura)
            **({"display_name": display_name} if display_name else {}),
            "iat": now,
            "exp": now + ttl,
            "jti": str(uuid.uuid4()),
        }

        token_value = _create_hs256_jwt(payload, secret)
        # La guía acepta token o cosmos_token; usamos cosmos_token
        params["cosmos_token"] = token_value

    return {
        "url": f"{callback}?{urlencode(params)}",
        "includes_cosmos_token": bool(token_value),
        "expires_in": ttl if token_value else None,
        "notetaker_base": base,
    }