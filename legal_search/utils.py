import os
import json
import uuid
import secrets
import time
import logging
import re
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import asyncio
import httpx
import pandas as pd
import docx2txt
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
from fastapi import Request, HTTPException, status
from jose import JWTError, jwt
from prometheus_client import Counter, Histogram
from sqlalchemy.orm import Session
import os
import socket
import struct
import time
import asyncio
import logging
from typing import Dict, Any, Optional, Set

from config.models import Conversation

# =============================================================================
#  CONFIGURACIÓN GENERAL
# =============================================================================

# Clave JWT y algoritmo: mismos que login / nlp / addUser
SECRET_KEY = os.getenv("COSMOS_SECRET_KEY", "secretkey123")
ALGORITHM = os.getenv("COSMOS_JWT_ALG", "HS256")

# Raíz de almacenamiento compartida por todos los servicios
STORAGE_ROOT = Path(os.getenv("COSMOS_STORAGE_ROOT", "./cosmos_data")).resolve()

# URLs de servicios internos (configurables para Docker o VS Code)
LLM_URL = os.getenv(
    "COSMOS_LLM_URL_RAG",
    "http://localhost:8081/generate_response_with_context_rag",
)
LLAMA_SERVER_URL = os.getenv(
    "COSMOS_LLM_URL_BASE",
    "http://localhost:8081/generate_response",
)
LLM_CONTEXT_CONVERSATION = os.getenv(
    "COSMOS_LLM_URL_CONV",
    "http://localhost:8081/generate_response_conversation",
)
LLM_AI_IMAGE = os.getenv(
    "COSMOS_LLM_URL_IMAGE",
    "http://localhost:8081/generate_response_ai_image",
)
LLM_LAST_DOCUMENT = os.getenv(
    "COSMOS_LLM_URL_LASTDOC",
    "http://localhost:8081/ask_last_document",
)

# OCR / NLP / TalkToDocument
OCR_APP_URL = os.getenv("COSMOS_OCR_URL", "http://localhost:6000")
NLP_APP_URL = os.getenv("COSMOS_NLP_URL", "http://localhost:5000")
TALK_TO_DOCUMENT_URL = os.getenv(
    "COSMOS_TALK_TO_DOC_URL", "http://localhost:5001/chat_document"
)
NLP_UPLOAD_FILE_URL = os.getenv(
    "COSMOS_NLP_UPLOAD_URL", "http://localhost:5000/upload_file"
)

# Roles (IDs numéricos en BD)
ROLE_SUPERVISOR = 1
ROLE_USER = 2

logger = logging.getLogger(__name__)

# =============================================================================
#  MÉTRICAS PROMETHEUS
# =============================================================================

REQUEST_COUNT = Counter(
    "http_requests_total",
    "Total de solicitudes HTTP",
    ["method", "endpoint"],
)

RESPONSE_TIME = Histogram(
    "http_response_time_seconds",
    "Tiempo de respuesta HTTP",
    ["method", "endpoint"],
)

HTTP_REQUEST_DURATION = Histogram(
    "http_request_duration_seconds",
    "Duración de las peticiones HTTP",
    ["method", "endpoint"],
)

# JWT verification
JWT_VERIFY_REQUESTS = Counter(
    "jwt_verify_requests_total",
    "Número de veces que se llamó a verify_token_from_cookie",
)
JWT_MISSING = Counter(
    "jwt_verify_missing_token_total",
    "Veces que no se encontró token en la cookie",
)
JWT_INVALID = Counter(
    "jwt_verify_invalid_token_total",
    "Veces que el token JWT fue inválido o expirado",
)
JWT_VALID = Counter(
    "jwt_verify_valid_token_total",
    "Veces que el token JWT fue verificado correctamente",
)
JWT_VERIFY_DURATION = Histogram(
    "jwt_verify_duration_seconds",
    "Tiempo que tarda la verificación de token JWT",
)

# Redis: sesión
REDIS_SESSION_REQUESTS = Counter(
    "redis_session_requests_total",
    "Total calls to get_session_from_redis",
)
REDIS_SESSION_HITS = Counter(
    "redis_session_hits_total",
    "Redis session cache hits",
)
REDIS_SESSION_MISSES = Counter(
    "redis_session_misses_total",
    "Redis session cache misses",
)
REDIS_SESSION_DURATION = Histogram(
    "redis_session_duration_seconds",
    "Duration of get_session_from_redis",
)

# Redis: conversaciones
CONV_GET_REQUESTS = Counter(
    "redis_conversation_get_requests_total",
    "Total calls to get_conversation_to_redis",
)
CONV_GET_HITS = Counter(
    "redis_conversation_get_hits_total",
    "Redis conversation cache hits",
)
CONV_GET_MISSES = Counter(
    "redis_conversation_get_misses_total",
    "Redis conversation cache misses",
)
CONV_GET_DURATION = Histogram(
    "redis_conversation_get_duration_seconds",
    "Duration of get_conversation_to_redis",
)

CONV_SAVE_REQUESTS = Counter(
    "redis_conversation_save_requests_total",
    "Total calls to save_conversation_to_redis",
)
CONV_SAVE_DURATION = Histogram(
    "redis_conversation_save_duration_seconds",
    "Duration of save_conversation_to_redis",
)

CONV_RESET_REQUESTS = Counter(
    "redis_conversation_reset_requests_total",
    "Total calls to reset_conversation_in_redis",
)
CONV_RESET_DURATION = Histogram(
    "redis_conversation_reset_duration_seconds",
    "Duration of reset_conversation_in_redis",
)

# Upload de ficheros efímeros
UPLOADFILE_REQUESTS = Counter(
    "uploadfile_requests_total",
    "Total calls to create_upload_file",
)
UPLOADFILE_ERRORS = Counter(
    "uploadfile_errors_total",
    "Total failures in create_upload_file",
)
UPLOADFILE_DURATION = Histogram(
    "uploadfile_duration_seconds",
    "Duration of create_upload_file handler",
    ["method", "endpoint"],
)

# LLM queries
llm_query_counter = Counter(
    "query_llm_requests_total",
    "Número total de peticiones a /query/llm",
)

llm_query_requests_by_status_total = Counter(
    "llm_query_requests_by_status_total",
    "Número total de peticiones a /query/llm clasificadas por status",
    ["status"],
)

llm_query_latency = Histogram(
    "llm_query_latency_seconds",
    "Latencia en segundos de /query/llm",
)

# Flujos por tipo (R/C)
query_flow_counter = Counter(
    "query_flow_requests_total",
    "Total requests per query flow",
    ["flow", "status"],
)
query_flow_latency = Histogram(
    "query_flow_request_latency_seconds",
    "Latency per query flow",
    ["flow"],
)

# Noticias / web
news_llm_query_counter = Counter(
    "news_llm_query_requests_total",
    "Total News LLM query requests",
)

news_llm_query_requests_by_status_total = Counter(
    "news_llm_query_requests_by_status_total",
    "News LLM query requests by status",
    ["status"],
)

news_llm_query_latency = Histogram(
    "news_llm_query_request_latency_seconds",
    "Latency of News LLM query requests",
)


# =============================================================================
#  EXTRACCIÓN DE TEXTO DE FICHEROS
# =============================================================================

def extract_text_from_pdf(file_path: str) -> str:
    return pdf_extract_text(file_path)


def extract_text_from_docx(file_path: str) -> str:
    return docx2txt.process(file_path)


def extract_text_from_doc(file_path: str) -> str:
    # Si en tu stack usas antiword u otra lib específica, cámbialo aquí.
    return docx2txt.process(file_path)


def extract_text_from_xlsx(file_path: str) -> str:
    df = pd.read_excel(file_path)
    records = df.astype(str).apply(lambda x: " ".join(x), axis=1).tolist()
    return "\n".join(records)


def extract_text_from_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    return "\n".join(df.astype(str).values.flatten())


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def allowed_file(filename: str) -> bool:
    """
    Comprueba si la extensión del fichero está soportada para subir
    como 'ephemeral file'.
    """
    return "." in filename and filename.rsplit(".", 1)[1].lower() in {
        "pdf",
        "pptx",
        "xlsx",
        "docx",
        "doc",
        "csv",
        "jpg",
        "jpeg",
        "png",
        "txt",
    }


# =============================================================================
#  LIMPIEZA DE RESPUESTAS DEL LLM
# =============================================================================

def clean_text(text: str) -> str:
    """
    Limpieza agresiva de la salida del LLM antes de devolverla al usuario:

    - Recorta espacios iniciales/finales.
    - Elimina fences ```...``` y etiquetas ```json.
    - Corta cualquier sección de razonamiento interno tipo "Thought:", "Razonamiento:", etc.
    - Quita disclaimers típicos.
    - Colapsa saltos de línea excesivos.
    """
    if not text:
        return ""

    t = str(text).strip()

    # 1) Quitar fences ```lang ... ```
    if t.startswith("```") and t.endswith("```"):
        t = re.sub(r"^```[a-zA-Z0-9_\-]*\s*\n?", "", t)
        t = re.sub(r"\n?```$", "", t)

    # 2) Cortar a partir de "Thought:" / "Razonamiento:" / "Explanation:"
    thought_patterns = [
        r"\n\s*Thought\s*:.*",
        r"\n\s*Razonamiento\s*:.*",
        r"\n\s*Explanation\s*:.*",
    ]
    for pat in thought_patterns:
        t = re.sub(pat, "", t, flags=re.IGNORECASE | re.DOTALL)

    # 3) Quitar cabeceras tipo "Respuesta:" al inicio
    t = re.sub(
        r"^(respuesta:|respuesta\s*|assistant:|asistente:)\s*",
        "",
        t,
        flags=re.IGNORECASE,
    )

    # 4) Eliminar disclaimers típicos
    disclaimer_patterns = [
        r"como modelo de lenguaje[^.\n]*\.",
        r"como ia(?: de)? lenguaje[^.\n]*\.",
    ]
    for pat in disclaimer_patterns:
        t = re.sub(pat, "", t, flags=re.IGNORECASE)

    # 5) Colapsar >2 saltos de línea
    t = re.sub(r"\n{3,}", "\n\n", t)

    return t.strip()


async def perform_rag_search(
    prompt: str,
    access_token: str,
    flow: str = "C",
    department_directory: Optional[str] = None,
) -> Dict[str, Any]:
    flow_norm = (flow or "C").upper()

    args: Dict[str, Any] = {
        "query": prompt,
        "flow": flow_norm,
        "access_token": access_token,
    }
    if department_directory:
        args["department_directory"] = department_directory

    # 1) Preferencia: MCP (streamable) si está disponible
    try:
        from mcp_client import MCPClient  # import local para evitar acoplar siempre
        mcp = MCPClient()
        rag_payload = await mcp.ainvoke_tool("rag_search_tool", args)
        if isinstance(rag_payload, dict):
            status = rag_payload.get("status") or "ok"
            return {
                "status": status,
                "results": rag_payload.get("results") or [],
                "aggregation": rag_payload.get("aggregation"),
                "used_department": rag_payload.get("used_department"),
                "cleaned_query": rag_payload.get("cleaned_query") or prompt,
                "message": rag_payload.get("message"),
            }
    except Exception as e:
        logger.warning("perform_rag_search: MCP fallback a REST por error: %s", e)

    # 2) Fallback: REST legacy
    mcp_base = (os.getenv("CREW_BASE_URL") or "http://127.0.0.1:8090/api/v1").rstrip("/")
    url = f"{mcp_base}/tools/rag_search_tool/invoke"
    payload = {"args": args}

    async with httpx.AsyncClient(timeout=180.0) as client:
        resp = await client.post(url, json=payload)
        resp.raise_for_status()
        data = resp.json() or {}

    rag_payload: Dict[str, Any] = data.get("result") or {}
    status = rag_payload.get("status") or "ok"

    return {
        "status": status,
        "results": rag_payload.get("results") or [],
        "aggregation": rag_payload.get("aggregation"),
        "used_department": rag_payload.get("used_department"),
        "cleaned_query": rag_payload.get("cleaned_query") or prompt,
        "message": rag_payload.get("message"),
    }



def perform_rag_search_sync(
    prompt: str,
    access_token: str,
    flow: str = "C",
    department_directory: Optional[str] = None,
) -> Dict[str, Any]:
    flow_norm = (flow or "C").upper()

    args: Dict[str, Any] = {
        "query": prompt,
        "flow": flow_norm,
        "access_token": access_token,
    }
    if department_directory:
        args["department_directory"] = department_directory

    # 1) Preferencia: MCP
    try:
        from mcp_client import MCPClient
        mcp = MCPClient()
        rag_payload = mcp.invoke_tool("rag_search_tool", args)
        if isinstance(rag_payload, dict):
            status = rag_payload.get("status") or "ok"
            return {
                "status": status,
                "results": rag_payload.get("results") or [],
                "aggregation": rag_payload.get("aggregation"),
                "used_department": rag_payload.get("used_department"),
                "cleaned_query": rag_payload.get("cleaned_query") or prompt,
                "message": rag_payload.get("message"),
            }
    except Exception as e:
        logger.warning("perform_rag_search_sync: MCP fallback a REST por error: %s", e)

    # 2) Fallback: REST legacy
    mcp_base = (os.getenv("CREW_BASE_URL") or "http://127.0.0.1:8090/api/v1").rstrip("/")
    url = f"{mcp_base}/tools/rag_search_tool/invoke"
    payload = {"args": args}

    with httpx.Client(timeout=180.0) as client:
        resp = client.post(url, json=payload)
    resp.raise_for_status()

    data = resp.json() or {}
    rag_payload: Dict[str, Any] = data.get("result") or {}
    status = rag_payload.get("status") or "ok"

    return {
        "status": status,
        "results": rag_payload.get("results") or [],
        "aggregation": rag_payload.get("aggregation"),
        "used_department": rag_payload.get("used_department"),
        "cleaned_query": rag_payload.get("cleaned_query") or prompt,
        "message": rag_payload.get("message"),
    }



# =============================================================================
#  JWT / AUTH BÁSICO
# =============================================================================
def verify_token(request: Request) -> Dict[str, Any]:
    """
    Verifica el JWT presente en la cookie 'access_token' o en el header Authorization.
    Lanza HTTPException(403) si falta o es inválido/expirado.
    Retorna el payload en caso de éxito.
    """
    JWT_VERIFY_REQUESTS.inc()
    start_time = time.time()

    access_token = request.cookies.get("access_token")
    token_str: Optional[str] = None

    if access_token:
        logger.info("✱ Extracting access token from session cookie...")
        token_str = access_token.replace("Bearer ", "").strip()
    else:
        auth = request.headers.get("Authorization")
        if auth and auth.lower().startswith("bearer "):
            logger.info("✱ Extracting access token from Authorization header...")
            token_str = auth[7:].strip()

    if not token_str:
        JWT_MISSING.inc()
        JWT_VERIFY_DURATION.observe(time.time() - start_time)
        logger.debug("✱ No token present in cookie/header.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="No autorizado: Token no presente.",
        )

    try:
        logger.info("✱ Decoding token to validate access...")
        payload = jwt.decode(token_str, SECRET_KEY, algorithms=[ALGORITHM])

        email = payload.get("sub")
        user_id = payload.get("user_id")
        if not email or user_id is None:
            JWT_INVALID.inc()
            JWT_VERIFY_DURATION.observe(time.time() - start_time)
            logger.debug("✱ Token payload missing 'sub' or 'user_id'.")
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Token no válido.",
            )

        JWT_VALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - start_time)
        logger.debug(
            "✱ Token verified successfully for user '%s' (id=%s).", email, user_id
        )
        return payload

    except JWTError:
        JWT_INVALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - start_time)
        logger.debug("✱ Token inválido o expirado.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Token no válido o expirado.",
        )


def verify_token_from_cookie(request: Request) -> str:
    """
    Verifica el JWT en la cookie 'access_token'.
    Retorna el email (campo 'sub') si es válido.
    Lanza HTTPException(403) en caso contrario.
    """
    JWT_VERIFY_REQUESTS.inc()
    t0 = time.time()

    access_token = request.cookies.get("access_token")
    if access_token is None:
        JWT_MISSING.inc()
        JWT_VERIFY_DURATION.observe(time.time() - t0)
        logger.debug("✱ Token no presente en la cookie.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="No autorizado: Token no presente.",
        )

    token_str = access_token.replace("Bearer ", "")
    try:
        payload = jwt.decode(token_str, SECRET_KEY, algorithms=[ALGORITHM])
        email = payload.get("sub")
        if not email:
            JWT_INVALID.inc()
            JWT_VERIFY_DURATION.observe(time.time() - t0)
            logger.debug("✱ Token no contiene 'sub'.")
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Token no válido.",
            )

        JWT_VALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - t0)
        return email

    except JWTError:
        JWT_INVALID.inc()
        JWT_VERIFY_DURATION.observe(time.time() - t0)
        logger.debug("✱ Token inválido o expirado.")
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Token no válido o expirado.",
        )


# =============================================================================
#  REDIS CLIENTS (inyectados desde main.py)
# =============================================================================

redis_client = None
redis_client_conversations = None


def init_redis_clients(core_client, conv_client) -> None:
    """
    Registra los clientes Redis que usará el microservicio de modelo de negocio.

    Debe llamarse desde main.py en el evento startup.
    """
    global redis_client, redis_client_conversations
    redis_client = core_client
    redis_client_conversations = conv_client
    logger.info(
        "init_redis_clients: Redis clients registrados (core db=0, conv db posiblemente 2)."
    )


# =============================================================================
#  CSRF (double-submit cookie)
# =============================================================================

def generate_csrf_token() -> str:
    """
    Genera un token CSRF aleatorio, URL-safe y con entropía suficiente.
    Compatible con lectura por el frontend (cookie no HttpOnly).
    """
    # 32 bytes ≈ 256 bits de entropía; se serializa base64-url
    return secrets.token_urlsafe(32)


def validate_csrf_double_submit(
    request: Request,
    *,
    cookie_name: str,
    header_name: str = "X-CSRFToken",
    error_detail: str = "CSRF token inválido o ausente.",
    allow_methods: Optional[Set[str]] = None,
) -> None:
    """
    Valida CSRF con el patrón "double-submit cookie":
      - La cookie `cookie_name` debe existir.
      - La cabecera `header_name` debe existir.
      - Ambos valores deben coincidir (comparación constante).

    Se ignoran métodos seguros (GET/HEAD/OPTIONS/TRACE) y los de `allow_methods`.
    """
    method = (request.method or "GET").upper()
    safe_methods = {"GET", "HEAD", "OPTIONS", "TRACE"}
    allow = set(allow_methods or set())

    # Para compatibilidad: si el validador se llama en endpoints idempotentes, no romper.
    if method in safe_methods or method in allow:
        return

    cookie_token = request.cookies.get(cookie_name)
    header_token = request.headers.get(header_name)

    # `request.headers` es case-insensitive, pero garantizamos str
    if not cookie_token or not header_token:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail=error_detail,
        )

    # Comparación en tiempo constante
    if not secrets.compare_digest(str(cookie_token), str(header_token)):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail=error_detail,
        )


# =============================================================================
#  HELPERS DE AUTH / COSTE
# =============================================================================

def extract_raw_bearer_token(request: Request) -> Optional[str]:
    """
    Extrae el JWT crudo (sin el prefijo 'Bearer ') desde:

      1) Cabecera Authorization: 'Bearer <jwt>'
      2) Cookie 'access_token': 'Bearer <jwt>'
    """
    auth_header = request.headers.get("Authorization", "")
    cookie_token = request.cookies.get("access_token")

    if auth_header.startswith("Bearer "):
        return auth_header.split(" ", 1)[1].strip()

    if cookie_token and cookie_token.startswith("Bearer "):
        return cookie_token.split(" ", 1)[1].strip()

    return None


def calculate_cost(total_tokens: int, price_per_1k: float = 0.0005) -> float:
    """
    Coste aproximado: price_per_1k por cada 1000 tokens.
    Se guarda como float en UsageLog.cost (Numeric(10,6)).
    """
    return round((total_tokens / 1000.0) * price_per_1k, 6)


# =============================================================================
#  CONVERSACIONES (BD)
# =============================================================================

def get_or_create_conversation(
    db: Session,
    user_id: int,
    conversation_id: Optional[int],
) -> Tuple[Conversation, bool]:
    """
    Devuelve una conversación existente del usuario o crea una nueva.

    return: (conversation, created)
      - created=True si se ha creado una nueva conversación.
      - created=False si se ha reutilizado una existente.
    """
    if conversation_id:
        convo = (
            db.query(Conversation)
            .filter(Conversation.id == conversation_id, Conversation.user_id == user_id)
            .first()
        )
        if convo:
            return convo, False

    convo = Conversation(
        user_id=user_id,
        conversation_text="",
        created_at=datetime.utcnow(),
    )
    db.add(convo)
    db.commit()
    db.refresh(convo)
    return convo, True


# =============================================================================
#  REDIS: SESIÓN + CONVERSACIÓN + ARCHIVOS EFÍMEROS
# =============================================================================

async def get_session_from_redis(user_id: int) -> Optional[Dict[str, Any]]:
    """
    Recupera la sesión del usuario desde Redis (clave session:{user_id}).
    """
    logger.info(
        "\n\n(utils - get_session_from_redis)  Llamada asíncrona a REDIS para obtener los datos de sesión..."
    )
    REDIS_SESSION_REQUESTS.inc()
    start = time.time()

    if redis_client is None:
        logger.error("get_session_from_redis: redis_client no inicializado.")
        REDIS_SESSION_MISSES.inc()
        REDIS_SESSION_DURATION.observe(time.time() - start)
        return None

    session = await redis_client.get(f"session:{user_id}")
    duration = time.time() - start
    REDIS_SESSION_DURATION.observe(duration)

    if session:
        REDIS_SESSION_HITS.inc()
        try:
            return json.loads(session)
        except json.JSONDecodeError:
            logger.warning(
                "get_session_from_redis: JSON corrupto para session:%s", user_id
            )
            REDIS_SESSION_MISSES.inc()
            return None

    REDIS_SESSION_MISSES.inc()
    return None


async def save_conversation_to_redis(
    user_id: int,
    expires_in: int = 480,
    conversation_entry: Optional[dict] = None,
) -> None:
    """
    Actualiza el historial de conversación en Redis (DB de conversaciones).
    Clave: conversation:{user_id}
    """
    logger.info(
        "\n\n(utils - save_conversation_to_redis)  Actualizando el historial de conversación en Redis para el usuario..."
    )
    CONV_SAVE_REQUESTS.inc()
    start = time.time()

    if redis_client_conversations is None:
        logger.error(
            "save_conversation_to_redis: redis_client_conversations no inicializado."
        )
        CONV_SAVE_DURATION.observe(time.time() - start)
        return

    conversation = await redis_client_conversations.get(f"conversation:{user_id}")
    if conversation:
        conversation_data = json.loads(conversation)
    else:
        conversation_data = {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_history": [],
        }

    if "conversation_history" not in conversation_data:
        conversation_data["conversation_history"] = []

    if conversation_entry and "conversation_id" in conversation_entry:
        conversation_data["current_conversation_id"] = conversation_entry[
            "conversation_id"
        ]

    conversation_data["conversation_history"].append(conversation_entry or {})

    logger.debug("save_conversation_to_redis: %s", conversation_data)

    await redis_client_conversations.set(
        f"conversation:{user_id}",
        json.dumps(conversation_data),
        ex=expires_in,
    )

    CONV_SAVE_DURATION.observe(time.time() - start)


async def add_ephemeral_file(
    user_id: int,
    filename: str,
    text: str,
    ttl: int = 3600,
    meta: Optional[Dict[str, Any]] = None,
) -> None:
    """
    Añade/actualiza la lista de archivos en vuelo del usuario en Redis:

    - Clave: ephemeral:{user_id}
    - Valor: lista JSON de objetos:
        { filename, text (truncado), uploaded_at, meta }
    """
    key = f"ephemeral:{user_id}"

    if redis_client is None:
        logger.error("add_ephemeral_file: redis_client no inicializado.")
        return

    existing = await redis_client.get(key)
    if existing:
        try:
            files = json.loads(existing)
        except json.JSONDecodeError:
            files = []
    else:
        files = []

    max_stored_chars = int(os.getenv("EPHEMERAL_MAX_STORED_CHARS", "12000"))

    files.append(
        {
            "filename": filename,
            "text": (text or "")[:max_stored_chars],
            "uploaded_at": datetime.utcnow().isoformat(),
            "meta": meta or {},
        }
    )

    await redis_client.set(key, json.dumps(files), ex=ttl)


async def get_ephemeral_files(user_id: int) -> List[dict]:
    """
    Devuelve la lista de archivos en vuelo del usuario.
    """
    key = f"ephemeral:{user_id}"

    if redis_client is None:
        logger.error("get_ephemeral_files: redis_client no inicializado.")
        return []

    existing = await redis_client.get(key)
    if not existing:
        return []
    try:
        return json.loads(existing)
    except json.JSONDecodeError:
        logger.warning("get_ephemeral_files: JSON corrupto para key=%s", key)
        return []


async def get_conversation_to_redis(user_id: int, expires_in: int = 480) -> dict:
    """
    Recupera la conversación asociada al usuario desde Redis (DB de conversaciones).
    Si no existe, se inicializa una nueva conversación con un ID único.
    """
    logger.info(
        "\n\n(utils - get_conversation_to_redis)  Recuperando la conversación en Redis para el usuario..."
    )
    CONV_GET_REQUESTS.inc()
    start = time.time()

    if redis_client_conversations is None:
        logger.error(
            "get_conversation_to_redis: redis_client_conversations no inicializado."
        )
        CONV_GET_MISSES.inc()
        CONV_GET_DURATION.observe(time.time() - start)
        return {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_history": [],
        }

    conversation = await redis_client_conversations.get(f"conversation:{user_id}")
    if conversation:
        CONV_GET_HITS.inc()
        conversation_data = json.loads(conversation)
        logger.debug(
            "\n\n(utils - get_conversation_to_redis)  Conversación recuperada: %s",
            conversation_data,
        )
    else:
        CONV_GET_MISSES.inc()
        conversation_data = {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_history": [],
        }
        logger.info(
            "\n\n(utils - get_conversation_to_redis)  No se encontró conversación. Inicializando nueva conversación..."
        )
        await redis_client_conversations.set(
            f"conversation:{user_id}",
            json.dumps(conversation_data),
            ex=expires_in,
        )

    CONV_GET_DURATION.observe(time.time() - start)
    return conversation_data


async def reset_conversation_in_redis(user_id: int, expires_in: int = 480) -> dict:
    """
    Resetea la conversación asociada al usuario en Redis (DB de conversaciones).
    """
    logger.info(
        "\n\n(utils - reset_conversation_in_redis)  Reseteando la conversación en Redis para el usuario..."
    )
    CONV_RESET_REQUESTS.inc()
    start = time.time()

    if redis_client_conversations is None:
        logger.error(
            "reset_conversation_in_redis: redis_client_conversations no inicializado."
        )
        CONV_RESET_DURATION.observe(time.time() - start)
        return {
            "conversation_id": str(uuid.uuid4()),
            "user_id": user_id,
            "conversation_history": [],
        }

    new_conversation_data = {
        "conversation_id": str(uuid.uuid4()),
        "user_id": user_id,
        "conversation_history": [],
    }
    await redis_client_conversations.set(
        f"conversation:{user_id}",
        json.dumps(new_conversation_data),
        ex=expires_in,
    )

    CONV_RESET_DURATION.observe(time.time() - start)
    return new_conversation_data


async def set_user_context(user_id: int, text: str, ttl: int = 3600) -> None:
    """
    Guarda el último contexto textual bajo la clave context:{user_id} en Redis.
    """
    if redis_client is None:
        logger.error("set_user_context: redis_client no inicializado.")
        return
    await redis_client.set(f"context:{user_id}", text, ex=ttl)


# =============================================================================
#  AUTH COMPUESTO: SESIÓN UNÍVOCA (JWT + REDIS)
# =============================================================================
async def get_current_auth_chatdoc(request: Request) -> Dict[str, Any]:
    """
    Dependencia de autenticación para microservicios tipo chat_document/web_search/legal_search.

    - Valida SIEMPRE el JWT (cookie 'access_token' o header Authorization) usando verify_token.
    - Recupera la sesión desde Redis (session:{user_id}).
    - Devuelve un dict compatible + añade `access_token` crudo (sin 'Bearer '),
      útil para propagar identidad a MCP (auditoría/ACL multiusuario).
    """
    # 1) Validar JWT
    token_payload = verify_token(request)
    user_id = token_payload.get("user_id")
    if not user_id:
        logger.warning("[auth] Token sin user_id.")
        raise HTTPException(status_code=401, detail="Token inválido (sin user_id).")

    # 2) Extraer JWT crudo (sin 'Bearer ')
    raw_jwt = extract_raw_bearer_token(request) or ""

    # 3) Recuperar sesión desde Redis
    try:
        session_data = await get_session_from_redis(user_id)
    except Exception as e:
        logger.warning("[auth] Error recuperando sesión en Redis para user_id=%s: %s", user_id, e)
        session_data = None

    if not isinstance(session_data, dict):
        session_data = {}

    return {
        "user_id": user_id,
        "role": token_payload.get("role"),
        "token_payload": token_payload,
        "session": session_data,
        # 🔹 NUEVO: passthrough multiusuario hacia MCP / tools
        "access_token": raw_jwt,
    }


# -------------------------------------------------
# CLAMAV
# -------------------------------------------------
async def scan_with_clamav(content: bytes, filename: Optional[str] = None) -> Dict[str, Any]:
    """
    Escanea un blob de bytes con ClamAV (clamd) vía INSTREAM.

    Devuelve un diccionario con:
      {
        "status": "OK" | "INFECTED" | "ERROR" | "SKIPPED",
        "virus_name": str | None,
        "raw_result": str | None,
        "error": str | None,
        "duration_s": float,
        "filename": str | None,
        "size_bytes": int
      }

    Política por defecto:
      - Si ClamAV responde 'FOUND' → status='INFECTED'
      - Si ClamAV responde 'OK'    → status='OK'
      - Cualquier otro resultado o excepción → status='ERROR'
    """
    start = time.time()

    # Fichero vacío → lo damos por OK
    if not content:
        return {
            "status": "OK",
            "virus_name": None,
            "raw_result": "EMPTY_FILE",
            "error": None,
            "duration_s": time.time() - start,
            "filename": filename,
            "size_bytes": 0,
        }

    # Host / puerto de ClamAV:
    # - En docker-compose, normalmente CLAMAV_HOST=clamav, CLAMAV_PORT=3310
    # - En host Windows con clamd publicado: CLAMAV_HOST=127.0.0.1
    host = os.getenv("CLAMAV_HOST", "127.0.0.1")
    port = int(os.getenv("CLAMAV_PORT", "3310"))
    timeout = float(os.getenv("CLAMAV_TIMEOUT", "10.0"))

    # Prefijo de comando: 'z' (NULL-terminated) o 'n' (newline-terminated).
    # ClamAV suele aceptar 'z' para TCP ('zINSTREAM\\0'), pero se puede ajustar
    # vía variable de entorno CLAMAV_CMD_PREFIX.
    cmd_prefix = os.getenv("CLAMAV_CMD_PREFIX", "z").lower()
    if cmd_prefix not in ("z", "n"):
        cmd_prefix = "z"

    loop = asyncio.get_running_loop()
    result = await loop.run_in_executor(
        None,
        _scan_with_clamav_sync,
        host,
        port,
        timeout,
        cmd_prefix,
        content,
        filename,
        start,
    )
    return result


def _scan_with_clamav_sync(
    host: str,
    port: int,
    timeout: float,
    cmd_prefix: str,
    content: bytes,
    filename: Optional[str],
    start: float,
) -> Dict[str, Any]:
    """
    Implementación síncrona del escaneo INSTREAM.

    - Intenta conectar primero al host configurado (CLAMAV_HOST).
    - Si falla la conexión, prueba hosts típicos en entorno Docker/Windows:
        * clamav
        * clamav_service
        * host.docker.internal
        * 127.0.0.1
        * localhost
    - Envía directamente el comando INSTREAM sin PING previo para evitar
      que algunos servidores cierren la conexión.
    """
    size_bytes = len(content)
    raw_result: Optional[str] = None
    virus_name: Optional[str] = None
    error: Optional[str] = None
    status = "ERROR"

    display_filename = filename or "<sin_nombre>"

    sock: Optional[socket.socket] = None
    last_exc: Optional[BaseException] = None

    # Fallbacks de host para soportar:
    #   - modelo de negocio en contenedor Docker
    #   - ClamAV en otro contenedor o en el host
    env_fallbacks = os.getenv("CLAMAV_HOST_FALLBACKS", "")
    extra_hosts: List[str] = [
        h.strip()
        for h in env_fallbacks.split(",")
        if h.strip()
    ]

    candidate_hosts: List[str] = [host] + extra_hosts
    for h in ("clamav", "clamav_service", "host.docker.internal", "127.0.0.1", "localhost"):
        if h not in candidate_hosts:
            candidate_hosts.append(h)

    host_used = host

    try:
        # 1) Conexión a alguno de los hosts candidatos
        for candidate in candidate_hosts:
            try:
                logger.info(
                    "[AV] Conectando a ClamAV en %s:%s para escanear '%s' (%d bytes)...",
                    candidate,
                    port,
                    display_filename,
                    size_bytes,
                )
                sock = socket.create_connection((candidate, port), timeout=timeout)
                host_used = candidate
                break
            except OSError as exc:
                last_exc = exc
                logger.warning(
                    "[AV] No se pudo conectar a ClamAV en %s:%s: %s",
                    candidate,
                    port,
                    exc,
                )

        if sock is None:
            # No se pudo conectar a ningún host candidato
            raise last_exc or RuntimeError(
                f"No se pudo conectar a ClamAV en {host}:{port}"
            )

        # 2) Preparar comando INSTREAM según prefijo z/n
        if cmd_prefix == "z":
            instream_cmd = b"zINSTREAM\0"
        else:  # "n"
            instream_cmd = b"nINSTREAM\n"

        chunk_size = int(os.getenv("CLAMAV_CHUNK_SIZE", "8192"))

        with sock:
            # 2) Protocolo INSTREAM
            # Prefijo INSTREAM + stream de chunks + chunk final de tamaño 0
            sock.sendall(instream_cmd)

            view = memoryview(content)
            offset = 0

            while offset < size_bytes:
                chunk = view[offset : offset + chunk_size]
                # Primero longitud (4 bytes BE) y luego los datos
                sock.sendall(struct.pack(">L", len(chunk)))
                sock.sendall(chunk)
                offset += len(chunk)

            # Fin de stream
            sock.sendall(struct.pack(">L", 0))

            # 3) Leer respuesta (ej: "stream: OK" o "...FOUND")
            chunks: List[bytes] = []
            while True:
                data = sock.recv(4096)
                if not data:
                    break
                chunks.append(data)

        raw_result = b"".join(chunks).decode(errors="replace").strip("\x00\r\n\t ")
        logger.info(
            "[AV] Resultado ClamAV para '%s' (host=%s): %s",
            display_filename,
            host_used,
            raw_result,
        )

        # 4) Interpretar resultado
        if "FOUND" in raw_result:
            status = "INFECTED"
            try:
                # "stream: <nombre_virus> FOUND"
                after_colon = raw_result.split(":", 1)[1].strip()
                virus_name = after_colon.rsplit(" ", 1)[0].strip()
            except Exception:
                virus_name = None
        elif "OK" in raw_result:
            status = "OK"
        else:
            status = "ERROR"
            error = f"Resultado inesperado de ClamAV: {raw_result}"

    except Exception as exc:
        # Cubre:
        #   - ConnectionRefused / timeout
        #   - ConnectionAborted (WinError 10053)
        #   - errores de protocolo
        error = str(exc)

        # Mensaje más explícito para el caso típico en Windows
        if isinstance(exc, ConnectionAbortedError):
            logger.error(
                "[AV] Conexión abortada al escanear '%s' con ClamAV (host=%s:%s): %s. "
                "En Windows esto suele indicar que el servidor cerró la conexión "
                "al recibir un comando inesperado o que un firewall/antivirus local "
                "interrumpió el flujo.",
                display_filename,
                host_used,
                port,
                exc,
                exc_info=True,
            )
        else:
            logger.error(
                "[AV] Error al escanear '%s' con ClamAV (host=%s:%s): %s",
                display_filename,
                host_used,
                port,
                exc,
                exc_info=True,
            )

        status = "ERROR"

    duration = time.time() - start

    return {
        "status": status,
        "virus_name": virus_name,
        "raw_result": raw_result,
        "error": error,
        "duration_s": duration,
        "filename": filename,
        "size_bytes": size_bytes,
    }

def _normalize_mcp_url(url: str) -> str:
    """
    Normaliza una URL base hacia el endpoint MCP streamable HTTP: .../mcp/

    Soporta entradas típicas:
      - http://host:8090
      - http://host:8090/api/v1
      - http://host:8090/mcp
      - http://host:8090/mcp/call_tool
    """
    u = (url or "").strip()
    if not u:
        return ""

    # Si llega base REST (/api/v1), la convertimos a /mcp
    if "/api/v1" in u:
        base = u.split("/api/v1", 1)[0].rstrip("/")
        u = base + "/mcp"

    # Si llega con subpaths MCP típicos
    u = re.sub(r"/mcp/(call_tool|list_tools|tools).*?$", "/mcp", u)

    # Asegurar sufijo /mcp/
    if u.endswith("/mcp"):
        u = u + "/"
    elif u.endswith("/mcp/"):
        pass
    else:
        u = u.rstrip("/") + "/mcp/"

    # Normalizar dobles slashes justo al final
    u = re.sub(r"/mcp/+$", "/mcp/", u)

    return u



